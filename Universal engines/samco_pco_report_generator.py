#!/usr/bin/env python3
"""
SAMCO-PCO Report Generator
Generates a 31-slide PowerPoint presentation from 30 CSV data files.
"""

import os
import csv
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
PRS = Presentation()
PRS.slide_width = Inches(13.333)
PRS.slide_height = Inches(7.5)

# Color palette
COLOR_TITLE = RGBColor(0, 51, 102)
COLOR_SUBTITLE = RGBColor(100, 100, 100)
COLOR_HEADER_BG = RGBColor(0, 51, 102)
COLOR_HEADER_TEXT = RGBColor(255, 255, 255)
COLOR_GREEN = RGBColor(34, 139, 34)
COLOR_RED = RGBColor(178, 34, 34)
COLOR_AMBER = RGBColor(218, 165, 32)
COLOR_BLUE = RGBColor(30, 144, 255)
COLOR_GRAY = RGBColor(128, 128, 128)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BLACK = RGBColor(0, 0, 0)


def read_csv(filename):
    path = os.path.join(OUTPUT_DIR, filename)
    with open(path, newline='', encoding='utf-8') as f:
        reader = csv.reader(f)
        headers = next(reader)
        rows = list(reader)
    return headers, rows


def add_title_slide(title, subtitle=""):
    blank = PRS.slide_layouts[6]
    slide = PRS.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PRS.slide_width, PRS.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_TITLE
    bg.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(1))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER
    return slide


def add_slide_header(slide, title, code, status=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PRS.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_HEADER_BG
    bar.line.fill.background()
    
    tbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9), Inches(0.6))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{code}  |  {title}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADER_TEXT
    
    if status:
        colors = {"PASS": COLOR_GREEN, "PASS WITH WARNINGS": COLOR_AMBER, 
                  "DRAFT": COLOR_GRAY, "FAIL": COLOR_RED}
        badge_color = colors.get(status, COLOR_BLUE)
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                        Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = badge_color
        badge.line.fill.background()
        btf = badge.text_frame
        bp = btf.paragraphs[0]
        bp.text = status
        bp.font.size = Pt(11)
        bp.font.bold = True
        bp.font.color.rgb = COLOR_WHITE
        bp.alignment = PP_ALIGN.CENTER


def add_kpi_box(slide, left, top, width, height, label, value, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    box.line.color.rgb = color
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = label
    p1.font.size = Pt(10)
    p1.font.color.rgb = COLOR_GRAY
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = str(value)
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = color
    p2.alignment = PP_ALIGN.CENTER


def add_data_table(slide, left, top, width, height, headers, rows, max_rows=8):
    num_rows = min(len(rows) + 1, max_rows + 1)
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLOR_WHITE
            paragraph.alignment = PP_ALIGN.CENTER
    for r_idx, row in enumerate(rows[:max_rows], 1):
        for c_idx, val in enumerate(row[:num_cols]):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.color.rgb = COLOR_BLACK
                paragraph.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
    return table


def add_text_box(slide, left, top, width, height, text, font_size=10, color=COLOR_BLACK, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    return box


# ==================== SLIDE GENERATORS ====================

def slide_01_tia():
    headers, rows = read_csv("PRO-TIA-01_Extension_of_Time.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Extension of Time / TIA", "PRO-TIA-01", "PASS WITH WARNINGS")
    total_eot = sum(int(r[9]) for r in rows if r[9].isdigit())
    pending = sum(1 for r in rows if r[8] == "Pending")
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Total EOT Days", total_eot, COLOR_BLUE)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Pending TIAs", pending, COLOR_AMBER)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Events", len(rows), COLOR_GREEN)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.8), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[2] for r in rows]
    chart_data.add_series("Duration", [int(r[5]) for r in rows])
    chart_data.add_series("Approved EOT", [int(r[9]) if r[9].isdigit() else 0 for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.5),
        "ANALYSIS NOTES:\n• 5 steel FIM delay events linked to late material delivery\n• 2 TIAs pending engineer review (>30 days)", 9, COLOR_GRAY)


def slide_02_delay():
    headers, rows = read_csv("PRO-DLY-02_Delay_Analysis.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Delay Analysis Report", "PRO-DLY-02", "PASS WITH WARNINGS")
    total_delay = sum(int(r[7]) for r in rows if r[7].isdigit())
    critical = sum(1 for r in rows if r[8] == "Critical")
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Total Delay Days", total_delay, COLOR_RED)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Critical Delays", critical, COLOR_RED)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Activities", len(rows), COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[2] for r in rows]
    chart_data.add_series("Delay Days", [int(r[7]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "ANALYSIS NOTES:\n• Foundation works delayed 28 days due to unforeseen ground conditions\n• Recovery plan: additional piling rig deployed", 9, COLOR_GRAY)


def slide_03_progress():
    headers, rows = read_csv("PRO-PRG-03_Project_Progress.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Project Progress Report", "PRO-PRG-03", "PASS")
    avg_planned = round(sum(float(r[2]) for r in rows) / len(rows), 1)
    avg_actual = round(sum(float(r[3]) for r in rows) / len(rows), 1)
    behind = sum(1 for r in rows if r[5] == "Behind")
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Avg Planned %", f"{avg_planned}%", COLOR_BLUE)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Avg Actual %", f"{avg_actual}%", COLOR_GREEN)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Behind Schedule", behind, COLOR_RED)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 8)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1] for r in rows]
    chart_data.add_series("Planned", [float(r[2]) for r in rows])
    chart_data.add_series("Actual", [float(r[3]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Structural works 3% ahead of plan\n• MEP coordination lagging due to design revisions", 9, COLOR_GRAY)


def slide_04_recovery():
    headers, rows = read_csv("PRO-REC-04_Recovery_Mitigation.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Recovery & Mitigation Report", "PRO-REC-04", "PASS")
    orig = sum(int(r[2]) for r in rows if r[2].isdigit())
    rec = sum(int(r[3]) for r in rows if r[3].isdigit())
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Original Impact", f"{orig} days", COLOR_RED)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Recovered", f"{rec} days", COLOR_GREEN)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Net Delay", f"{orig-rec} days", COLOR_AMBER)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1] for r in rows]
    chart_data.add_series("Original", [int(r[2]) for r in rows])
    chart_data.add_series("Recovered", [int(r[3]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Parallel trades strategy recovering 15 days on superstructure\n• Overtime budget: $120K approved", 9, COLOR_GRAY)


def slide_05_variation():
    headers, rows = read_csv("PRO-VAR-05_Variation_Report.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Variation Report", "PRO-VAR-05", "PASS")
    total_est = sum(float(r[5]) for r in rows)
    total_app = sum(float(r[6]) for r in rows)
    pending = sum(1 for r in rows if r[4] == "Pending")
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Total Est. Value", f"${total_est/1e6:.2f}M", COLOR_BLUE)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Approved Value", f"${total_app/1e6:.2f}M", COLOR_GREEN)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Pending", pending, COLOR_AMBER)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1][:20] for r in rows]
    chart_data.add_series("Estimated", [float(r[5])/1000 for r in rows])
    chart_data.add_series("Approved", [float(r[6])/1000 for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Road widening variation ($420K) approved by owner\n• Landscape upgrade under review - value engineering ongoing", 9, COLOR_GRAY)


def slide_06_ipc():
    headers, rows = read_csv("PRO-IPC-06_Interim_Payment.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Interim Payment Certificate", "PRO-IPC-06", "PASS")
    total_cert = float(rows[-1][7])
    total_net = float(rows[-1][9])
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Total Certified", f"${total_cert/1e6:.2f}M", COLOR_BLUE)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Net Paid", f"${total_net/1e6:.2f}M", COLOR_GREEN)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Retention", f"${float(rows[-1][8])/1e6:.2f}M", COLOR_AMBER)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[0] for r in rows]
    chart_data.add_series("This Period", [float(r[5])/1000 for r in rows])
    chart_data.add_series("Cumulative", [float(r[7])/1000 for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• IPC-12 certified with minor deductions for incomplete NCR closures\n• Payment cycle: 28 days from certification", 9, COLOR_GRAY)


def slide_07_resource():
    headers, rows = read_csv("PRO-RES-07_Resource_Manpower.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Resource / Manpower Report", "PRO-RES-07", "PASS")
    total_head = sum(int(r[2]) for r in rows)
    actual_head = sum(int(r[3]) for r in rows)
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Planned Headcount", total_head, COLOR_BLUE)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Actual Headcount", actual_head, COLOR_GREEN)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Variance", actual_head - total_head, COLOR_AMBER if actual_head >= total_head else COLOR_RED)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 8)
    chart_data = CategoryChartData()
    chart_data.categories = [r[0] for r in rows]
    chart_data.add_series("Planned", [int(r[1]) for r in rows])
    chart_data.add_series("Actual", [int(r[2]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• MEP technician shortage (-12) being addressed via agency recruitment\n• Steel fixer surplus (+8) reallocated to secondary structure", 9, COLOR_GRAY)


def slide_08_equipment():
    headers, rows = read_csv("PRO-EQP-08_Equipment_Plant.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Equipment & Plant Report", "PRO-EQP-08", "PASS")
    total_idle = sum(float(r[6]) for r in rows)
    avg_util = round(sum(float(r[10]) for r in rows) / len(rows), 1)
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Total Idle Hours", f"{total_idle:.0f}h", COLOR_AMBER)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Avg Utilization", f"{avg_util}%", COLOR_GREEN if avg_util > 75 else COLOR_AMBER)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Fleet Size", len(rows), COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 8)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1][:15] for r in rows]
    chart_data.add_series("Planned Hrs", [float(r[4]) for r in rows])
    chart_data.add_series("Actual Hrs", [float(r[5]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Tower crane TC-1 availability at 94% - maintenance on schedule\n• Concrete pump idle hours elevated due to pour sequence changes", 9, COLOR_GRAY)


def slide_09_cost():
    headers, rows = read_csv("PRO-CST-09_Cost_Control.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Cost Control Report", "PRO-CST-09", "PASS")
    total_bud = sum(float(r[2]) for r in rows)
    total_fc = sum(float(r[5]) for r in rows)
    var = total_fc - total_bud
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Total Budget", f"${total_bud/1e6:.2f}M", COLOR_BLUE)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Forecast Final", f"${total_fc/1e6:.2f}M", COLOR_GREEN if var <= 0 else COLOR_RED)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Variance", f"${var/1e6:.2f}M", COLOR_GREEN if var <= 0 else COLOR_RED)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1] for r in rows]
    chart_data.add_series("Budget", [float(r[2])/1000 for r in rows])
    chart_data.add_series("Forecast", [float(r[5])/1000 for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Materials forecast 4.2% over budget due to steel price escalation\n• Contingency drawdown at 35% - within approved limits", 9, COLOR_GRAY)


def slide_10_evm():
    headers, rows = read_csv("PRO-EVM-10_Earned_Value.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Earned Value Management", "PRO-EVM-10", "PASS")
    avg_spi = round(sum(float(r[8]) for r in rows) / len(rows), 3)
    avg_cpi = round(sum(float(r[9]) for r in rows) / len(rows), 3)
    total_vac = sum(float(r[12]) for r in rows)
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Avg SPI", avg_spi, COLOR_GREEN if avg_spi >= 1 else COLOR_RED)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Avg CPI", avg_cpi, COLOR_GREEN if avg_cpi >= 1 else COLOR_RED)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "VAC", f"${total_vac/1e6:.2f}M", COLOR_GREEN if total_vac >= 0 else COLOR_RED)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 5)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1] for r in rows]
    chart_data.add_series("SPI", [float(r[8]) for r in rows])
    chart_data.add_series("CPI", [float(r[9]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• SPI trending upward - schedule recovery measures effective\n• CPI stable at 0.97; cost underrun expected in MEP package", 9, COLOR_GRAY)


def slide_11_cashflow():
    headers, rows = read_csv("PRO-CFS-11_Cash_Flow.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Cash Flow Report", "PRO-CFS-11", "PASS")
    cum_net = float(rows[-1][7])
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Cumulative Net", f"${cum_net/1e6:.2f}M", COLOR_GREEN if cum_net > 0 else COLOR_RED)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Months Tracked", len(rows), COLOR_BLUE)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Avg Monthly Net", f"${cum_net/len(rows)/1e6:.2f}M", COLOR_GREEN)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[0] for r in rows]
    chart_data.add_series("Planned Net", [float(r[3])/1000 for r in rows])
    chart_data.add_series("Actual Net", [float(r[6])/1000 for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Cash flow positive for 14 consecutive months\n• Q3 outflow spike expected due to facade package mobilization", 9, COLOR_GRAY)


def slide_12_scurve():
    headers, rows = read_csv("PRO-SCV-12_S_Curve.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "S-Curve Performance", "PRO-SCV-12", "PASS")
    final_actual = float(rows[-1][2])
    final_forecast = float(rows[-1][3])
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Current Actual", f"{final_actual}%", COLOR_GREEN)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Forecast Final", f"{final_forecast}%", COLOR_BLUE)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Variance", f"{final_forecast - final_actual:.1f}%", COLOR_AMBER)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[0] for r in rows]
    chart_data.add_series("Planned", [float(r[1]) for r in rows])
    chart_data.add_series("Actual", [float(r[2]) for r in rows])
    chart_data.add_series("Forecast", [float(r[4]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• S-curve tracking within 3% of planned trajectory\n• Recovery measures expected to close gap by Q4", 9, COLOR_GRAY)


def slide_13_procurement():
    headers, rows = read_csv("PRO-PRC-13_Procurement_Status.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Procurement Status", "PRO-PRC-13", "PASS WITH WARNINGS")
    total_val = sum(float(r[8]) for r in rows)
    delayed = sum(1 for r in rows if r[7] == "Delayed")
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Total Value", f"${total_val/1e6:.2f}M", COLOR_BLUE)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Delayed", delayed, COLOR_RED)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Packages", len(rows), COLOR_GREEN)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1][:15] for r in rows]
    chart_data.add_series("Value ($K)", [float(r[8])/1000 for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Elevator package delayed 21 days - manufacturer capacity constraint\n• MEP package on track for Q3 award", 9, COLOR_GRAY)


def slide_14_material():
    headers, rows = read_csv("PRO-MAT-14_Material_Status.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Material Status", "PRO-MAT-14", "PASS WITH WARNINGS")
    short = sum(1 for r in rows if r[8] in ["Short", "Critical Short"])
    total_req = sum(float(r[3]) for r in rows)
    total_del = sum(float(r[4]) for r in rows)
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Shortfalls", short, COLOR_RED)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Delivered %", f"{total_del/total_req*100:.1f}%", COLOR_GREEN)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Materials", len(rows), COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1][:15] for r in rows]
    chart_data.add_series("Required", [float(r[3]) for r in rows])
    chart_data.add_series("Delivered", [float(r[4]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Glass panels critical short - alternative supplier engaged\n• Rebar delivery ahead of schedule (+200MT buffer)", 9, COLOR_GRAY)


def slide_15_labour():
    headers, rows = read_csv("PRO-LAH-15_Labour_Hours.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Labour & Hours Report", "PRO-LAH-15", "PASS")
    total_hours = sum(float(r[4]) for r in rows)
    total_cost = sum(float(r[7]) for r in rows)
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Total Hours", f"{total_hours/1000:.1f}K", COLOR_BLUE)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Total Cost", f"${total_cost/1e6:.2f}M", COLOR_GREEN)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Trades", len(set(r[1] for r in rows)), COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    trade_data = {}
    for r in rows:
        trade = r[1]
        trade_data[trade] = trade_data.get(trade, 0) + float(r[4])
    chart_data = CategoryChartData()
    chart_data.categories = list(trade_data.keys())
    chart_data.add_series("Hours", list(trade_data.values()))
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Overtime hours within approved limits (12% of regular)\n• Productivity factor averaging 0.94 across all trades", 9, COLOR_GRAY)


def slide_16_cpm():
    headers, rows = read_csv("PRO-CPM-16_Critical_Path.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Critical Path Method", "PRO-CPM-16", "PASS WITH WARNINGS")
    critical = sum(1 for r in rows if r[11] == "Yes")
    total_dur = sum(int(r[2]) for r in rows)
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Critical Activities", critical, COLOR_RED)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Total Duration", f"{total_dur} days", COLOR_BLUE)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Float < 10d", sum(1 for r in rows if r[10].isdigit() and int(r[10]) < 10), COLOR_AMBER)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1] for r in rows]
    chart_data.add_series("Duration", [int(r[2]) for r in rows])
    chart_data.add_series("Total Float", [int(r[7]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Critical path runs through structure and MEP rough-in\n• Commissioning activity has 15-day float - buffer for handover", 9, COLOR_GRAY)


def slide_17_float():
    headers, rows = read_csv("PRO-FLT-17_Float_Analysis.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Float Analysis", "PRO-FLT-17", "PASS WITH WARNINGS")
    high_risk = sum(1 for r in rows if r[6] == "High")
    avg_tf = round(sum(int(r[2]) for r in rows) / len(rows), 1)
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "High Risk", high_risk, COLOR_RED)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Avg Total Float", f"{avg_tf} days", COLOR_BLUE)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Activities", len(rows), COLOR_GREEN)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1] for r in rows]
    chart_data.add_series("Total Float", [int(r[2]) for r in rows])
    chart_data.add_series("Free Float", [int(r[3]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "ANALYSIS NOTES:\n• 3 activities with < 10 days total float require monitoring\n• Steel erection free float consumed by predecessor delay", 9, COLOR_GRAY)


def slide_18_milestone():
    headers, rows = read_csv("PRO-MIL-18_Milestone_Status.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Milestone Status", "PRO-MIL-18", "PASS")
    achieved = sum(1 for r in rows if r[5] == "Achieved")
    at_risk = sum(1 for r in rows if r[5] == "At Risk")
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Achieved", achieved, COLOR_GREEN)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "At Risk", at_risk, COLOR_RED)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Total", len(rows), COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1] for r in rows]
    chart_data.add_series("Variance (days)", [int(r[6]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Structure topped out achieved 5 days ahead of forecast\n• Practical completion at risk - facade delays may impact Gate 3", 9, COLOR_GRAY)


def slide_19_risk():
    headers, rows = read_csv("PRO-RSK-19_Risk_Report.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Project Risk Report", "PRO-RSK-19", "PASS WITH WARNINGS")
    high = sum(1 for r in rows if r[5] == "High")
    open_risks = sum(1 for r in rows if r[8] == "Open")
    total_cont = sum(float(r[9]) for r in rows)
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "High Risks", high, COLOR_RED)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Open", open_risks, COLOR_AMBER)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Contingency", f"${total_cont/1e6:.2f}M", COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1][:15] for r in rows]
    chart_data.add_series("Probability", [int(r[3]) for r in rows])
    chart_data.add_series("Impact", [int(r[4]) * 20 for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "ANALYSIS NOTES:\n• Steel price volatility (High) - hedge contract executed for Q3-Q4\n• Labor shortage risk escalated to owner - visa processing delays", 9, COLOR_GRAY)


def slide_20_rfi():
    headers, rows = read_csv("PRO-RFI-20_RFI_Submittal.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "RFI / Submittal Performance", "PRO-RFI-20", "PASS")
    overdue = sum(1 for r in rows if r[6] == "Overdue")
    avg_response = round(sum(int(r[4]) for r in rows) / len(rows), 1)
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Overdue", overdue, COLOR_RED)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Avg Response", f"{avg_response} days", COLOR_GREEN if avg_response <= 10 else COLOR_AMBER)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Total RFIs", len(rows), COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[0] for r in rows]
    chart_data.add_series("Response Days", [int(r[4]) for r in rows])
    chart_data.add_series("Target", [int(r[5]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• RFI response rate at 92% within target (7-14 days)\n• 2 overdue RFIs escalated to engineer for priority review", 9, COLOR_GRAY)


def slide_21_qaqc():
    headers, rows = read_csv("PRO-QAQC-21_QAQC_Performance.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "QA/QC Performance", "PRO-QAQC-21", "PASS")
    pass_rate = sum(1 for r in rows if r[4] == "Pass") / len(rows) * 100
    total_nc = sum(int(r[5]) for r in rows)
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Pass Rate", f"{pass_rate:.0f}%", COLOR_GREEN)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Total NCRs", total_nc, COLOR_AMBER)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Inspections", len(rows), COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1][:15] for r in rows]
    chart_data.add_series("NCR Count", [int(r[5]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Concrete slab L3 inspection passed with zero NCRs\n• Fireproofing thickness NCR closed after re-inspection", 9, COLOR_GRAY)


def slide_22_productivity():
    headers, rows = read_csv("PRO-PROD-22_Productivity.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Productivity Performance", "PRO-PROD-22", "PASS")
    avg_ratio = round(sum(float(r[4]) for r in rows) / len(rows), 3)
    below_target = sum(1 for r in rows if float(r[4]) < 1.0)
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Avg Ratio", avg_ratio, COLOR_GREEN if avg_ratio >= 1 else COLOR_AMBER)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Below Target", below_target, COLOR_RED)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Activities", len(rows), COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[0] for r in rows]
    chart_data.add_series("Planned", [float(r[2]) for r in rows])
    chart_data.add_series("Actual", [float(r[3]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• MEP rough-in productivity at 1.08 - above target\n• Plastering below target (0.82) - skill gap training scheduled", 9, COLOR_GRAY)


def slide_23_executive():
    headers, rows = read_csv("PRO-EXE-23_Executive_Report.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Executive Project Report", "PRO-EXE-23", "PASS WITH WARNINGS")
    behind = sum(1 for r in rows if r[6] == "Behind")
    caution = sum(1 for r in rows if r[6] == "Caution")
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Behind", behind, COLOR_RED)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Caution", caution, COLOR_AMBER)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "On Track", sum(1 for r in rows if r[6] == "On Track"), COLOR_GREEN)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 8)
    add_text_box(slide, Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.5),
        "EXECUTIVE SUMMARY:\n\n"
        "• Project 68.5% complete (2.5% behind plan)\n"
        "• CPI 0.97 - slight cost overrun contained\n"
        "• Zero LTIs - safety performance exemplary\n"
        "• Recovery plan on track for Q4\n"
        "• 12 open NCRs - closure rate improving\n"
        "• Manpower gap of 38 being addressed\n"
        "• Equipment utilization rising to 78.5%\n"
        "• Forecast completion: Q1 2027", 11, COLOR_BLACK)
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Executive dashboard updated weekly - next review 20-Aug-2026", 9, COLOR_GRAY)


def slide_24_baseline():
    headers, rows = read_csv("PRO-BLC-24_Baseline_vs_Current.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Baseline vs Current Schedule", "PRO-BLC-24", "PASS")
    total_shift = sum(int(r[8]) for r in rows if r[8].isdigit())
    severe = sum(1 for r in rows if r[9] == "Severe")
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Total Shift", f"{total_shift} days", COLOR_AMBER)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Severe Impact", severe, COLOR_RED)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Activities", len(rows), COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1] for r in rows]
    chart_data.add_series("Start Variance", [int(r[6]) for r in rows])
    chart_data.add_series("Finish Variance", [int(r[7]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Foundation start shifted +12 days due to access road delay\n• Finishes finish variance -8 days (ahead) due to early mobilization", 9, COLOR_GRAY)


def slide_25_forecast():
    headers, rows = read_csv("PRO-FCF-25_Forecast_Completion.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Forecast Completion", "PRO-FCF-25", "PASS")
    likely_date = [r[1] for r in rows if r[0] == "Most Likely"][0]
    risk_adj_date = [r[1] for r in rows if r[0] == "Risk-Adjusted"][0]
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Most Likely", likely_date, COLOR_BLUE)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Risk-Adjusted", risk_adj_date, COLOR_AMBER)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Scenarios", len(rows), COLOR_GREEN)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 4)
    chart_data = CategoryChartData()
    chart_data.categories = [r[0] for r in rows]
    chart_data.add_series("Duration", [int(r[2]) for r in rows])
    chart_data.add_series("Probability", [int(r[3]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Risk-adjusted forecast shows 35-day extension vs baseline\n• Contingency reserve adequate for pessimistic scenario", 9, COLOR_GRAY)


def slide_26_ml():
    headers, rows = read_csv("PRO-ML-26_ML_Project_Controls.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "ML Project Controls", "PRO-ML-26", "DRAFT")
    avg_acc = round(sum(float(r[3]) for r in rows) / len(rows), 1)
    prod = sum(1 for r in rows if r[8] == "Production")
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Avg Accuracy", f"{avg_acc}%", COLOR_BLUE)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "In Production", prod, COLOR_GREEN)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Models", len(rows), COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 5)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1] for r in rows]
    chart_data.add_series("Accuracy", [float(r[3]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Cost Forecaster at 92.3% accuracy - highest performing model\n• Productivity Estimator in development - target Q4 deployment", 9, COLOR_GRAY)


def slide_27_contract():
    headers, rows = read_csv("PRO-CAD-27_Contract_Admin.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Contract Administration", "PRO-CAD-27", "PASS")
    open_items = sum(1 for r in rows if r[3] == "Open")
    total_val = sum(float(r[6]) for r in rows)
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Open Items", open_items, COLOR_RED)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Total Value", f"${total_val/1e6:.2f}M", COLOR_BLUE)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Items", len(rows), COLOR_GREEN)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1][:15] for r in rows]
    chart_data.add_series("Value ($K)", [float(r[6])/1000 for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Final account negotiation ongoing - $1.8M disputed amount\n• Retention release application submitted for Gate 3 milestone", 9, COLOR_GRAY)


def slide_28_change():
    headers, rows = read_csv("PRO-CHG-28_Change_Control.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Change Control Register", "PRO-CHG-28", "PASS")
    approved = sum(1 for r in rows if r[5] == "Approved")
    pending = sum(1 for r in rows if r[5] == "Pending")
    total_cost = sum(float(r[6]) for r in rows if r[6])
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Approved", approved, COLOR_GREEN)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "Pending", pending, COLOR_AMBER)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Total Cost", f"${total_cost/1e6:.2f}M", COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1][:15] for r in rows]
    chart_data.add_series("Cost Impact", [float(r[6])/1000 for r in rows])
    chart_data.add_series("Time Impact", [int(r[7]) for r in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Change control board meets bi-weekly - next session 18-Aug\n• All changes tracked against Clause 13.1 - zero unauthorized works", 9, COLOR_GRAY)


def slide_29_document():
    headers, rows = read_csv("PRO-DOC-29_Document_Control.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Document Control", "PRO-DOC-29", "PASS")
    approved = sum(1 for r in rows if r[4] == "Approved")
    for_approval = sum(1 for r in rows if r[4] == "For Approval")
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Approved", approved, COLOR_GREEN)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "For Approval", for_approval, COLOR_AMBER)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Documents", len(rows), COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[2] for r in rows]
    chart_data.add_series("Count", [1 for _ in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Document turnover rate: 94% within 14 days\n• Fire Strategy Report approved - issued for construction", 9, COLOR_GRAY)


def slide_30_interface():
    headers, rows = read_csv("PRO-IFC-30_Interface_Management.csv")
    slide = PRS.slides.add_slide(PRS.slide_layouts[6])
    add_slide_header(slide, "Interface Management", "PRO-IFC-30", "PASS")
    open_if = sum(1 for r in rows if r[4] == "Open")
    high_pri = sum(1 for r in rows if r[5] == "High")
    add_kpi_box(slide, Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.7), "Open", open_if, COLOR_RED)
    add_kpi_box(slide, Inches(3.0), Inches(1.1), Inches(2.5), Inches(0.7), "High Priority", high_pri, COLOR_RED)
    add_kpi_box(slide, Inches(5.7), Inches(1.1), Inches(2.5), Inches(0.7), "Interfaces", len(rows), COLOR_BLUE)
    add_data_table(slide, Inches(0.3), Inches(2.0), Inches(8.5), Inches(4.5), headers, rows, 6)
    chart_data = CategoryChartData()
    chart_data.categories = [r[1][:15] for r in rows]
    chart_data.add_series("Count", [1 for _ in rows])
    x, y, cx, cy = Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    add_text_box(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.6),
        "NOTES:\n• Structural-MEP coordination L5 resolved - BIM clash detection complete\n• Road-drainage alignment under review - survey data pending", 9, COLOR_GRAY)


# ==================== MAIN EXECUTION ====================

if __name__ == "__main__":
    print("Generating SAMCO-PCO Report Presentation...")
    
    # Title slide
    add_title_slide(
        "SAMCO-PCO Project Controls Report",
        "30 Template Slides | Structural Discipline | August 2026"
    )
    
    # Generate all 30 report slides
    slide_01_tia()
    slide_02_delay()
    slide_03_progress()
    slide_04_recovery()
    slide_05_variation()
    slide_06_ipc()
    slide_07_resource()
    slide_08_equipment()
    slide_09_cost()
    slide_10_evm()
    slide_11_cashflow()
    slide_12_scurve()
    slide_13_procurement()
    slide_14_material()
    slide_15_labour()
    slide_16_cpm()
    slide_17_float()
    slide_18_milestone()
    slide_19_risk()
    slide_20_rfi()
    slide_21_qaqc()
    slide_22_productivity()
    slide_23_executive()
    slide_24_baseline()
    slide_25_forecast()
    slide_26_ml()
    slide_27_contract()
    slide_28_change()
    slide_29_document()
    slide_30_interface()
    
    # Save
    output_path = os.path.join(OUTPUT_DIR, "SAMCO_PCO_31_Slides_Report.pptx")
    PRS.save(output_path)
    print(f"Done! Saved to: {output_path}")
    print(f"Total slides: {len(PRS.slides)}")