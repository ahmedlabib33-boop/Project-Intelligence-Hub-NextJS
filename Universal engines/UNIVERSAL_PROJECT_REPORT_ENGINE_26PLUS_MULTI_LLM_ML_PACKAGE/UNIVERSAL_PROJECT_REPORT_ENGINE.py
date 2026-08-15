#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
UNIVERSAL PROJECT REPORT ENGINE

Generic, company-neutral source code for AI chats, coding agents, and local use.
It reads PDF, DOCX, XLSX/XLSM, Primavera P6 XER, TXT, CSV, JSON, and ZIP inputs,
detects the report type, builds a normalized project model, and generates:

- Fully editable PowerPoint made from native text and shapes
- PNG-based PowerPoint
- Editable SVG pages
- High-resolution PNG pages
- Full-bleed A4 landscape PDF
- Project intelligence JSON
- Validation report
- ZIP package

No company name, logo, event value, project value, or project-specific evidence is
embedded. Uploaded evidence determines content. Branding is optional by config.
"""
from __future__ import annotations

import argparse, csv, io, json, re, shutil, tempfile, textwrap, zipfile, hashlib, platform
from difflib import SequenceMatcher
from dataclasses import dataclass, field, asdict
from datetime import datetime, date
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence, Tuple
from xml.sax.saxutils import escape as xml_escape

VERSION = "1.3.1"
W, H = 1920, 1080
SUPPORTED = {".pdf", ".docx", ".xlsx", ".xlsm", ".xer", ".txt", ".md", ".csv", ".json", ".zip"}
ALLOWED_REPORT_TYPES = {"auto","eot","tia","delay","progress","recovery","variation","hybrid"}
THEME = {
    "navy":"#06294F", "navy2":"#0B3B68", "gold":"#F8B915", "white":"#FFFFFF",
    "ink":"#102A43", "muted":"#64748B", "bg":"#F3F7FB", "border":"#BED0E2",
    "green":"#168E47", "red":"#D5242B", "purple":"#7041AD", "orange":"#D76A00",
    "teal":"#008F99", "gray":"#65717D", "pale_green":"#EAF7EF", "pale_red":"#FDEBED",
    "pale_purple":"#F3ECFA", "pale_orange":"#FFF2E3", "pale_gray":"#EEF2F5"
}
DEFAULT_BRANDING = {
    "company_name":"Project Controls", "prepared_by":"Project Controls Team",
    "confidentiality":"Confidential", "footer_left":None, "footer_right":None
}

@dataclass
class Source:
    path: str
    kind: str
    title: str
    text: str = ""
    tables: List[Dict[str, Any]] = field(default_factory=list)
    meta: Dict[str, Any] = field(default_factory=dict)
    warnings: List[str] = field(default_factory=list)

@dataclass
class Schedule:
    path: str
    role: str
    pair_key: str
    project: str = ""
    data_date: Optional[str] = None
    finish: Optional[str] = None
    finish_name: Optional[str] = None
    task_ids: List[str] = field(default_factory=list)
    rel_ids: List[str] = field(default_factory=list)
    tasks: List[Dict[str,str]] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)

@dataclass
class Comparison:
    name: str
    before: str
    after: str
    before_finish: Optional[str]
    after_finish: Optional[str]
    movement_days: Optional[int]
    inserted_tasks: int
    inserted_relationships: int
    confidence: str

@dataclass
class Event:
    event_id: str
    title: str
    description: str = ""
    cause: str = "Requires verification"
    effect: str = "Requires verification"
    movement_days: Optional[int] = None
    treatment: str = "Assessment required"
    evidence: List[str] = field(default_factory=list)
    confidence: str = "Medium"

@dataclass
class Model:
    project_name: str
    report_type: str
    confidence: float
    title: str
    data_date: Optional[str]
    period: str
    sources: List[Source]
    schedules: List[Schedule]
    comparisons: List[Comparison]
    events: List[Event]
    metrics: Dict[str, Any]
    progress: Dict[str, Any]
    milestones: List[Dict[str, Any]]
    constraints: List[str]
    risks: List[str]
    actions: List[str]
    conclusions: List[str]
    warnings: List[str]
    generated_at: str = field(default_factory=lambda: datetime.now().isoformat(timespec="seconds"))

@dataclass
class E:
    kind: str
    x: float
    y: float
    w: float = 0
    h: float = 0
    text: str = ""
    fill: str = "#FFFFFF"
    stroke: str = "#000000"
    sw: float = 0
    radius: float = 0
    fs: float = 24
    color: str = "#102A43"
    bold: bool = False
    align: str = "left"
    valign: str = "top"
    opacity: float = 1.0

@dataclass
class Slide:
    title: str
    elements: List[E] = field(default_factory=list)
    def add(self, e:E): self.elements.append(e)

# --------------------------- Helpers ---------------------------------------

def s(v:Any)->str:
    if v is None: return ""
    if isinstance(v,(datetime,date)): return v.isoformat()
    return str(v)

def clean(t:str)->str:
    return re.sub(r"\n{3,}","\n\n",re.sub(r"[ \t]+"," ",(t or "").replace("\x00"," "))).strip()

def slug(t:str)->str:
    return re.sub(r"[^A-Za-z0-9]+","_",t or "").strip("_") or "item"

def parse_date(v:Any)->Optional[datetime]:
    if isinstance(v,datetime): return v
    if isinstance(v,date): return datetime(v.year,v.month,v.day)
    t=s(v).strip()
    for f in ["%Y-%m-%d %H:%M:%S","%Y-%m-%d","%d-%b-%Y %H:%M:%S","%d-%b-%Y","%d/%m/%Y","%d-%m-%Y","%m/%d/%Y"]:
        try: return datetime.strptime(t,f)
        except: pass
    try: return datetime.fromisoformat(t.replace("Z","+00:00"))
    except: return None

def dtext(v:Any)->Optional[str]:
    x=parse_date(v); return x.strftime("%d-%b-%Y") if x else (s(v).strip() or None)

def days(a:Any,b:Any)->Optional[int]:
    x,y=parse_date(a),parse_date(b)
    return (y.date()-x.date()).days if x and y else None

def num(v:Any)->Optional[float]:
    m=re.search(r"[-+]?\d+(?:\.\d+)?",s(v).replace(",",""))
    return float(m.group()) if m else None

def pct(v:Any)->Optional[float]:
    n=num(v)
    if n is None: return None
    if 0<=n<=1 and "%" not in s(v): n*=100
    return n

def load_json(p:Path)->Any: return json.loads(p.read_text(encoding="utf-8-sig"))
def save_json(p:Path,v:Any): p.parent.mkdir(parents=True,exist_ok=True); p.write_text(json.dumps(v,indent=2,ensure_ascii=False,default=str),encoding="utf-8")

MAX_ZIP_MEMBER_BYTES = 250 * 1024 * 1024
MAX_ZIP_TOTAL_BYTES = 2 * 1024 * 1024 * 1024


def _safe_extract_zip(archive:Path,dest:Path)->List[Path]:
    """Extract supported files without allowing path traversal or oversized members."""
    extracted=[]; total=0; root=dest.resolve()
    with zipfile.ZipFile(archive) as z:
        for info in z.infolist():
            if info.is_dir(): continue
            total += max(0,info.file_size)
            if info.file_size>MAX_ZIP_MEMBER_BYTES:
                raise ValueError(f"ZIP member is too large: {info.filename}")
            if total>MAX_ZIP_TOTAL_BYTES:
                raise ValueError(f"ZIP expands beyond the permitted size: {archive.name}")
            target=(dest/info.filename).resolve()
            if root!=target and root not in target.parents:
                raise ValueError(f"Unsafe ZIP path rejected: {info.filename}")
            if target.suffix.lower() not in SUPPORTED-{".zip"}: continue
            target.parent.mkdir(parents=True,exist_ok=True)
            with z.open(info) as src,target.open('wb') as dst: shutil.copyfileobj(src,dst)
            extracted.append(target)
    return extracted


def collect_inputs(items:Sequence[str|Path], temp:Path)->List[Path]:
    if not items: raise ValueError("At least one input file or folder is required.")
    out=[]
    for item in items:
        p=Path(item).expanduser().resolve()
        if not p.exists(): raise FileNotFoundError(p)
        if p.is_dir(): out += [x for x in sorted(p.rglob("*")) if x.is_file() and x.suffix.lower() in SUPPORTED]
        elif p.suffix.lower()==".zip":
            dest=temp/f"unzipped_{slug(p.stem)}"; dest.mkdir(parents=True,exist_ok=True)
            out += _safe_extract_zip(p,dest)
        elif p.suffix.lower() in SUPPORTED: out.append(p)
    seen=set(); result=[]
    for path in out:
        k=str(path.resolve()).casefold()
        if k not in seen: seen.add(k); result.append(path)
    if not result: raise ValueError("No supported evidence files were found in the supplied inputs.")
    return result

def branding(config:Optional[str|Path]=None)->Dict[str,Any]:
    b=dict(DEFAULT_BRANDING)
    if config: b.update(load_json(Path(config)))
    b["footer_left"]=b.get("footer_left") or f"{b['company_name']} | {b['confidentiality']}"
    b["footer_right"]=b.get("footer_right") or f"Prepared by | {b['prepared_by']}"
    return b

# --------------------------- Readers ---------------------------------------

def read_pdf(p:Path, ocr:bool=False)->Source:
    try: from pypdf import PdfReader
    except ImportError: raise RuntimeError("Install pypdf")
    r=PdfReader(str(p)); parts=[]; warns=[];tables=[]
    for i,page in enumerate(r.pages):
        try: parts.append(page.extract_text() or "")
        except Exception as e: warns.append(f"Page {i+1} extraction failed: {e}")
    text=clean("\n\n".join(parts))
    # Optional structured table extraction. Failure does not block narrative extraction.
    try:
        import pdfplumber
        with pdfplumber.open(str(p)) as pdf:
            for pi,page in enumerate(pdf.pages[:200],1):
                for ti,table in enumerate(page.extract_tables() or [],1):
                    rows=[[s(cell).strip() for cell in row] for row in table if row]
                    if not rows:continue
                    h,rr=infer_table(rows)
                    if h and rr:tables.append({"name":f"Page {pi} Table {ti}","headers":h,"rows":rr[:2000],"row_count":len(rr)})
    except ImportError:
        warns.append("pdfplumber is not installed; PDF table extraction was skipped.")
    except Exception as e:
        warns.append(f"PDF table extraction was skipped after an error: {e}")
    if len(text)/max(1,len(r.pages))<80: warns.append("Low-text/scanned PDF detected. Use AI vision or reviewed OCR and pass a context JSON.")
    return Source(str(p),"pdf",p.stem,text,tables=tables,meta={"pages":len(r.pages),"characters":len(text),"tables":len(tables)},warnings=warns)

def read_docx(p:Path)->Source:
    try: from docx import Document
    except ImportError: raise RuntimeError("Install python-docx")
    d=Document(str(p)); parts=[]; tables=[]
    for q in d.paragraphs:
        if q.text.strip(): parts.append(q.text.strip())
    for i,t in enumerate(d.tables,1):
        rows=[[c.text.strip() for c in r.cells] for r in t.rows]
        h=rows[0] if rows else []; rr=rows[1:] if len(rows)>1 else []
        tables.append({"name":f"Table {i}","headers":h,"rows":rr})
        parts += [" | ".join(x) for x in rows[:100]]
    return Source(str(p),"docx",p.stem,clean("\n".join(parts)),tables,{"paragraphs":len(d.paragraphs),"tables":len(tables)})

def infer_table(rows:List[List[Any]])->Tuple[List[str],List[List[Any]]]:
    if not rows:return [],[]
    best=0; score=-1
    for i,row in enumerate(rows[:20]):
        vals=[x for x in row if s(x).strip()]
        sc=len(vals)+2*sum(bool(re.search(r"[A-Za-z]",s(x))) for x in vals)
        if len(vals)>=2 and sc>score: best,score=i,sc
    h=[s(x).strip() or f"Column_{i+1}" for i,x in enumerate(rows[best])]; w=len(h); out=[]
    for row in rows[best+1:]:
        x=list(row[:w])+[""]*max(0,w-len(row))
        if any(s(v).strip() for v in x): out.append(x)
    return h,out

def read_excel(p:Path)->Source:
    try: from openpyxl import load_workbook
    except ImportError: raise RuntimeError("Install openpyxl")
    wb=load_workbook(str(p),read_only=True,data_only=True); tables=[]; parts=[]
    for ws in wb.worksheets:
        rows=[]
        for i,row in enumerate(ws.iter_rows(values_only=True),1):
            vals=[s(x).strip() for x in row]
            while vals and vals[-1]=="": vals.pop()
            if any(vals): rows.append(vals)
            if i>=5000: break
        h,rr=infer_table(rows); tables.append({"name":ws.title,"headers":h,"rows":rr,"row_count":len(rr)})
        parts += [f"Worksheet: {ws.title}"," | ".join(h)] + [" | ".join(map(s,x)) for x in rr[:200]]
    return Source(str(p),"excel",p.stem,clean("\n".join(parts)),tables,{"sheets":wb.sheetnames})

def read_text(p:Path)->Source:
    if p.suffix.lower()==".json": return Source(str(p),"json",p.stem,json.dumps(load_json(p),indent=2,ensure_ascii=False,default=str))
    if p.suffix.lower()==".csv":
        rows=[]
        with p.open("r",encoding="utf-8-sig",errors="replace") as f:
            for i,row in enumerate(csv.reader(f)): rows.append(row); 
        h=rows[0] if rows else []; rr=rows[1:]
        return Source(str(p),"csv",p.stem,clean("\n".join(" | ".join(x) for x in rows[:500])),[{"name":p.stem,"headers":h,"rows":rr}])
    return Source(str(p),p.suffix[1:],p.stem,clean(p.read_text(encoding="utf-8-sig",errors="replace")))

def _xer_name_tokens(name:str)->List[str]:
    text=re.sub(r"[_\-.()]+"," ",name.lower())
    return [x for x in re.findall(r"[a-z0-9]+",text) if x not in {
        "after","before","befor","post","pre","impacted","unimpacted","fragnet","fargent","network","schedule","xer","dd"
    }]


def _xer_role(name:str)->str:
    tokens=set(_xer_name_tokens(name)+re.findall(r"[a-z]+",re.sub(r"[_\-.()]+"," ",name.lower())))
    if tokens & {"after","post","impacted"}: return "after"
    if tokens & {"before","befor","pre","unimpacted"}: return "before"
    return "other"


def _task_finish_date(row:Dict[str,str])->Optional[datetime]:
    for fld in ["early_end_date","remain_early_end_date","target_end_date","expect_end_date","act_end_date","end_date"]:
        x=parse_date(row.get(fld))
        if x:return x
    return None


def _finish_candidate_score(row:Dict[str,str])->int:
    name=(row.get("task_name") or "").lower(); code=(row.get("task_code") or row.get("task_id") or "").lower(); typ=(row.get("task_type") or "").lower()
    score=0
    if re.search(r"\bproject\s+(finish|completion|complete)\b",name): score+=160
    if re.search(r"\b(overall|contract)\s+(finish|completion)\b",name): score+=120
    if "project" in name and any(k in name for k in ["finish","completion","complete"]): score+=100
    if "mile" in typ: score+=35
    if typ in {"tt_finish","finish milestone"}: score+=40
    if any(k in code for k in ["proj_finish","project_finish","completion"]): score+=40
    if any(k in name for k in ["section","building","floor","ground works"]): score-=25
    return score


def parse_xer(p:Path)->Tuple[Source,Schedule]:
    tables={}; table=None; fields=[]
    with p.open("r",encoding="utf-8-sig",errors="replace") as f:
        for raw in f:
            q=raw.rstrip("\r\n").split("\t"); mark=q[0]
            if mark=="%T": table=q[1] if len(q)>1 else "UNKNOWN"; tables.setdefault(table,[]); fields=[]
            elif mark=="%F" and table: fields=q[1:]
            elif mark=="%R" and table and fields:
                vals=q[1:]+[""]*max(0,len(fields)-len(q[1:])); tables[table].append(dict(zip(fields,vals)))
            elif mark=="%E": break
    role=_xer_role(p.stem); tokens=_xer_name_tokens(p.stem); key=" ".join(tokens).strip() or slug(p.stem).lower()
    proj=(tables.get("PROJECT") or [{}])[0]; project=proj.get("proj_short_name") or proj.get("proj_name") or p.stem
    dd=dtext(proj.get("last_recalc_date") or proj.get("data_date"))
    tasks=tables.get("TASK",[]); rel=tables.get("TASKPRED",[])
    candidates=[]
    for row in tasks:
        dt=_task_finish_date(row); score=_finish_candidate_score(row)
        if dt and score>0:candidates.append((score,dt,row))
    best=None; best_date=None
    if candidates:
        _,best_date,best=max(candidates,key=lambda x:(x[0],x[1]));
    if not best_date:
        # PROJECT scheduled finish is a controlled fallback, not a substitute for a verified finish milestone.
        for fld in ["scd_end_date","plan_end_date","target_end_date","last_recalc_date"]:
            x=parse_date(proj.get(fld))
            if x:best_date=x;break
    warn=[]
    if not best_date: warn.append("Project Finish could not be identified from XER TASK or PROJECT data.")
    elif not best: warn.append("Project Finish used PROJECT-level fallback because no reliable finish milestone was identified.")
    task_ids=[x.get("task_id","") for x in tasks if x.get("task_id")]
    rel_ids=["|".join([x.get("task_id",""),x.get("pred_task_id",""),x.get("pred_type",""),x.get("lag_hr_cnt","")]) for x in rel]
    sched=Schedule(str(p),role,key,project,dd,dtext(best_date),best.get("task_name") if best else "PROJECT scheduled finish",task_ids,rel_ids,tasks,warn)
    src=Source(str(p),"xer",p.stem,f"Project: {project}\nRole: {role}\nData date: {dd}\nProject finish: {sched.finish}\nTasks: {len(tasks)}\nRelationships: {len(rel)}",meta={"tasks":len(tasks),"relationships":len(rel),"finish":sched.finish,"pair_key":key,"role":role},warnings=warn)
    return src,sched

def read_source(p:Path)->Tuple[Source,Optional[Schedule]]:
    e=p.suffix.lower()
    if e==".pdf": return read_pdf(p),None
    if e==".docx": return read_docx(p),None
    if e in {".xlsx",".xlsm"}: return read_excel(p),None
    if e==".xer": return parse_xer(p)
    return read_text(p),None

# --------------------------- Analysis --------------------------------------

KEYWORDS={
 "eot":{"extension of time":8,"time impact analysis":8,"eot":6,"fragnet":6,"concurrency":5,"before fragnet":5,"after fragnet":5},
 "delay":{"delay analysis":8,"windows analysis":7,"critical path":4,"delay event":4,"schedule impact":4,"float":2},
 "progress":{"progress report":8,"planned progress":6,"actual progress":6,"weekly report":5,"monthly report":5,"look ahead":5,"spi":4},
 "recovery":{"recovery plan":8,"acceleration":5,"mitigation plan":5,"catch-up":5},
 "variation":{"variation":7,"change order":6,"additional works":5,"scope change":5}
}

def classify(src:List[Source], sch:List[Schedule])->Tuple[str,float,Dict[str,float]]:
    corpus="\n".join(x.title+"\n"+x.text for x in src).lower(); scores={k:0.0 for k in KEYWORDS}
    for k,dd in KEYWORDS.items():
        for phrase,w in dd.items(): scores[k]+=w*min(corpus.count(phrase),5)
    if any(x.role=="before" for x in sch) and any(x.role=="after" for x in sch): scores["eot"]+=15; scores["delay"]+=10
    ordered=sorted(scores.items(),key=lambda x:x[1],reverse=True); best,bs=ordered[0]; second=ordered[1][1]
    if bs<=0:return "hybrid",.25,scores
    if second and bs/second<1.2 and {best,ordered[1][0]}!={"eot","delay"}: best="hybrid"
    return best,round(min(.99,.45+(bs-second)/max(bs,1)*.45),2),scores

def _pair_similarity(a:Schedule,b:Schedule)->float:
    name=SequenceMatcher(None,a.pair_key,b.pair_key).ratio()
    project=SequenceMatcher(None,(a.project or "").lower(),(b.project or "").lower()).ratio()
    return .8*name+.2*project


def compare(schedules:List[Schedule])->List[Comparison]:
    before=[x for x in schedules if x.role=="before"]; after=[x for x in schedules if x.role=="after"]
    out=[]; used=set()
    for b in sorted(before,key=lambda x:x.path):
        candidates=[(i,a,_pair_similarity(b,a)) for i,a in enumerate(after) if i not in used]
        if not candidates:continue
        i,a,sim=max(candidates,key=lambda x:x[2])
        if sim<.42:continue
        used.add(i); key=b.pair_key if b.pair_key==a.pair_key else SequenceMatcher(None,b.pair_key,a.pair_key).get_matching_blocks() and b.pair_key
        out.append(Comparison(re.sub(r"\s+"," ",key).title() or "Schedule Comparison",b.path,a.path,b.finish,a.finish,days(b.finish,a.finish),len(set(a.task_ids)-set(b.task_ids)),len(set(a.rel_ids)-set(b.rel_ids)),"High" if sim>=.75 and b.finish and a.finish else "Medium"))
    return out

def project_name(src:List[Source],sch:List[Schedule])->str:
    for x in sch:
        if x.project and x.project!=Path(x.path).stem:return x.project
    corpus="\n".join(x.text for x in src)
    for pat in [r"(?im)^\s*project(?: name)?\s*[:\-]\s*(.{4,120})$",r"(?im)^\s*contract\s*[:\-]\s*(.{4,120})$"]:
        m=re.search(pat,corpus)
        if m:return m.group(1).strip()[:140]
    return next((x.title for x in src if len(x.title)>8),"Project Report")

def detect_data_date(corpus:str,sch:List[Schedule])->Optional[str]:
    for pat in [r"(?im)data date\s*[:\-]\s*([^\n|]{6,30})",r"(?im)status date\s*[:\-]\s*([^\n|]{6,30})",r"(?im)cut[- ]off date\s*[:\-]\s*([^\n|]{6,30})"]:
        m=re.search(pat,corpus)
        if m:return dtext(m.group(1))
    vals=[parse_date(x.data_date) for x in sch if x.data_date]; vals=[x for x in vals if x]
    return dtext(max(vals)) if vals else None

def progress_values(src:List[Source])->Dict[str,Any]:
    found={"planned":[],"actual":[],"schedule":[]}
    for x in src:
        for t in x.tables:
            h=[s(v).lower() for v in t.get("headers",[])]; idx={}
            for i,v in enumerate(h):
                vv=re.sub(r"[^a-z0-9%]+"," ",v)
                if any(k in vv for k in ["planned","plan","baseline","target"]) and any(k in vv for k in ["%","progress","complete","completion"]):idx["planned"]=i
                if any(k in vv for k in ["actual","achieved","physical"]) and any(k in vv for k in ["%","progress","complete","completion"]):idx["actual"]=i
                if any(k in vv for k in ["schedule","time"]) and any(k in vv for k in ["%","progress","complete","completion"]):idx["schedule"]=i
            for row in t.get("rows",[]):
                for k,i in idx.items():
                    if i<len(row):
                        n=pct(row[i])
                        if n is not None and -10<=n<=110:found[k].append(n)
    r={k:round(v[-1],2) for k,v in found.items() if v}
    # Also read explicit percentages from narrative PDF/Word/text evidence.
    corpus="\n".join(x.text for x in src)
    text_patterns={
        "planned":[r"planned\s+(?:progress|complete|completion)\s*[:=\-]?\s*(\d+(?:\.\d+)?)\s*%"],
        "actual":[r"actual\s+(?:progress|complete|completion)\s*[:=\-]?\s*(\d+(?:\.\d+)?)\s*%"],
        "schedule":[r"schedule\s+(?:progress|complete|completion)\s*[:=\-]?\s*(\d+(?:\.\d+)?)\s*%"],
    }
    for key,patterns in text_patterns.items():
        if key in r: continue
        for pattern in patterns:
            hits=re.findall(pattern,corpus,re.I)
            if hits:
                r[key]=round(float(hits[-1]),2)
                break
    if "planned" in r and "actual" in r:r["variance"]=round(r["actual"]-r["planned"],2);r["spi_proxy"]=round(r["actual"]/r["planned"],3) if r["planned"] else None
    return r

def find_metric(corpus:str,patterns:List[str])->Optional[float]:
    for p in patterns:
        m=re.findall(p,corpus,re.I)
        if m:return float(m[-1])
    return None

def milestones(src:List[Source],sch:List[Schedule])->List[Dict[str,Any]]:
    out=[]
    for x in sch:
        for row in x.tasks:
            name=row.get("task_name",""); typ=row.get("task_type","").lower()
            if "mile" not in typ and not re.search(r"milestone|completion|finish",name,re.I):continue
            dt=next((dtext(row.get(f)) for f in ["early_end_date","target_end_date","expect_end_date","act_end_date","end_date"] if row.get(f)),None)
            out.append({"milestone":name or row.get("task_code","Milestone"),"date":dt,"status":"From XER","source":Path(x.path).name})
    seen=set();r=[]
    for x in out:
        k=(x["milestone"]+s(x["date"])).lower()
        if k not in seen:seen.add(k);r.append(x)
    return r[:50]

def list_items(corpus:str,words:List[str])->List[str]:
    out=[]
    for line in corpus.splitlines():
        q=line.strip(" \t•-–—")
        if 12<=len(q)<=250 and len(q.split())>=3 and any(w in q.lower() for w in words):out.append(q)
        if len(out)>=20:break
    return list(dict.fromkeys(out))

def events(src:List[Source], comps:List[Comparison], report_type:str)->List[Event]:
    out=[]
    for i,c in enumerate(comps,1):
        move=c.movement_days
        out.append(Event(f"EV{i:02d}",c.name,f"Before/after XER comparison between {Path(c.before).name} and {Path(c.after).name}.","Delay or change represented by the supplied schedule pair.",f"Project Finish moved by {move} calendar days." if move is not None else "Project Finish movement could not be calculated.",move,"Measured Project Finish movement" if move and move>0 else "No measured Project Finish movement" if move==0 else "Movement requires validation",[f"Before XER: {Path(c.before).name}",f"After XER: {Path(c.after).name}",f"Inserted activities: {c.inserted_tasks}",f"Inserted relationships: {c.inserted_relationships}"],c.confidence))
    corpus="\n".join(x.text for x in src)
    for rid,title in re.findall(r"(?im)^\s*((?:EV|Event|Delay Event)[-_ ]?\d{1,3})\s*[:\-–—]\s*(.{4,160})$",corpus):
        title=title.strip();number=re.search(r"\d+",rid);number=int(number.group()) if number else None
        matched=None
        for existing in out:
            exnum=re.search(r"\d+",existing.event_id)
            if number is not None and exnum and int(exnum.group())==number:matched=existing;break
            if title.lower() in existing.title.lower() or existing.title.lower() in title.lower():matched=existing;break
        if matched:
            if re.fullmatch(r"event\s*0*\d+",matched.title.strip(),re.I) or len(matched.title.strip())<9:matched.title=title
            if matched.description.startswith("Before/after XER comparison"):
                matched.description=f"{title}. "+matched.description
            ref=f"Narrative event reference: {rid} - {title}"
            if ref not in matched.evidence:matched.evidence.append(ref)
            continue
        out.append(Event(slug(rid).upper(),title,"Event extracted from submitted narrative evidence.","Requires source verification","Requires schedule linkage verification",None,"Extracted event; further analysis required",[],"Medium"))
        if len(out)>=12:break
    return out[:12]

def validate_context(context:Optional[Dict[str,Any]])->Tuple[Dict[str,Any],List[str]]:
    if context is None:return {},[]
    if not isinstance(context,dict):raise TypeError("Context must be a JSON object / Python dictionary.")
    out=dict(context);warnings=[]
    if out.get("report_type") and out["report_type"] not in ALLOWED_REPORT_TYPES-{"auto","tia"}:
        warnings.append(f"Unsupported context report_type '{out['report_type']}' was ignored.");out.pop("report_type",None)
    for section in ["metrics","progress"]:
        if section in out and not isinstance(out[section],dict):
            warnings.append(f"Context field '{section}' must be an object and was ignored.");out.pop(section,None)
    for section in ["milestones","constraints","risks","actions","conclusions","warnings","events"]:
        if section in out and not isinstance(out[section],list):
            warnings.append(f"Context field '{section}' must be a list and was ignored.");out.pop(section,None)
    numeric_fields={"gross_impact_days","overlap_days","net_eot_days"}
    if isinstance(out.get("metrics"),dict):
        for key in list(out["metrics"]):
            if key in numeric_fields:
                value=num(out["metrics"][key])
                if value is None or value<0:
                    warnings.append(f"Context metric '{key}' is invalid and was ignored.");out["metrics"].pop(key,None)
                else:out["metrics"][key]=value
    return out,warnings


def file_sha256(path:Path)->str:
    h=hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda:f.read(1024*1024),b""):h.update(chunk)
    return h.hexdigest()


def source_inventory(sources:List[Source])->List[Dict[str,Any]]:
    rows=[]
    for i,src in enumerate(sources,1):
        path=Path(src.path);rows.append({
            "source_id":src.meta.get("source_id",f"SRC-{i:03d}"),"file_name":path.name,"kind":src.kind,
            "size_bytes":path.stat().st_size if path.exists() else None,"sha256":file_sha256(path) if path.exists() else None,
            "title":src.title,"warnings":src.warnings,"metadata":src.meta
        })
    return rows


def context_template()->Dict[str,Any]:
    return {
        "project_name":"Project Name","report_type":"auto","data_date":"DD-Mmm-YYYY","period":"Reporting Period",
        "metrics":{"gross_impact_days":None,"overlap_days":None,"net_eot_days":None},
        "progress":{"planned":None,"actual":None},
        "events":[{"event_id":"EV01","title":"Verified event title","description":"Verified narrative","cause":"Verified cause","effect":"Verified schedule effect","movement_days":None,"treatment":"Assessment required","confidence":"Medium","evidence":["Source reference"]}],
        "milestones":[],"constraints":[],"risks":[],"actions":[],"conclusions":[],"warnings":[]
    }


def build_model(src:List[Source],sch:List[Schedule],rtype:str="auto",context:Optional[Dict[str,Any]]=None)->Model:
    context,context_warnings=validate_context(context)
    detected,conf,scores=classify(src,sch); rt=detected if rtype=="auto" else ("eot" if rtype=="tia" else rtype)
    if rt not in {"eot","delay","progress","recovery","variation","hybrid"}:rt="hybrid"
    corpus="\n".join(x.text for x in src); comps=compare(sch); ev=events(src,comps,rt); prog=progress_values(src)
    metrics={}
    pos=[x.movement_days for x in ev if x.movement_days is not None and x.movement_days>0]
    if pos:
        metrics["gross_impact_days"]=sum(pos)
        metrics["gross_impact_basis"]="Sum of positive movements from paired before/after schedules; validate that the comparisons are independent before contractual use."
    metrics["overlap_days"]=find_metric(corpus,[r"(?:concurrent|concurrency|overlap)[^\d]{0,25}(\d+)\s*(?:calendar\s*)?days"])
    metrics["net_eot_days"]=find_metric(corpus,[r"net\s+(?:eot|extension|impact)[^\d]{0,20}(\d+)\s*(?:calendar\s*)?days"])
    metrics["net_eot_basis"]="Explicitly extracted from source evidence" if metrics["net_eot_days"] is not None else "Not established automatically"
    af=[x.after_finish for x in comps if x.after_finish]; bf=[x.before_finish for x in comps if x.before_finish]
    if af:metrics["latest_finish"]=max(af,key=lambda x:parse_date(x) or datetime.min)
    if bf:metrics["base_finish"]=min(bf,key=lambda x:parse_date(x) or datetime.max)
    metrics["classification_scores"]=scores
    warnings=list(context_warnings)
    for x in src:warnings+=x.warnings
    unmatched=[Path(x.path).name for x in sch if x.role=="other"]
    if unmatched:warnings.append("XER role could not be inferred for: "+", ".join(unmatched[:8]))
    for x in sch:warnings+=x.warnings
    if rt in {"eot","delay"} and not comps:warnings.append("No paired before/after XER comparison identified.")
    if rt=="progress" and not prog:warnings.append("Planned and actual progress were not identified confidently.")
    conclusions=[]
    if rt in {"eot","delay"}:
        if pos:conclusions.append(f"The supplied paired schedules produce {sum(pos)} gross calendar days of measured Project Finish movement before concurrency treatment.")
        else:conclusions.append("No reliable positive Project Finish movement was established automatically.")
        if metrics.get("overlap_days") is not None and metrics.get("net_eot_days") is not None:
            ov=int(metrics["overlap_days"]); net=int(metrics["net_eot_days"])
            conclusions.append(f"The supplied evidence states {ov} days of overlap and a net position of {net} days; both require activity-level and Project Finish validation.")
        elif metrics.get("overlap_days") is not None:
            conclusions.append(f"The supplied evidence identifies {int(metrics['overlap_days'])} days of overlap, but the engine does not infer a net EOT automatically without a verified calculation basis.")
        else:conclusions.append("A net EOT position should not be stated until concurrency/overlap is validated.")
    elif rt=="progress":
        if "planned" in prog and "actual" in prog:
            conclusions.append(f"Actual progress is {abs(prog['variance']):.2f} percentage points {'ahead of' if prog['variance']>=0 else 'behind'} plan.")
        else:conclusions.append("Planned and actual progress values require confirmation.")
    else:conclusions.append("The engine consolidated the submitted evidence into a structured project-controls view.")
    conclusions.append("Responsibility, entitlement, contractual completion, and final schedule conclusions remain subject to native source verification and formal determination where applicable.")
    model=Model(project_name(src,sch),rt,conf,{"eot":"Extension of Time / Time Impact Analysis Report","delay":"Delay Analysis Report","progress":"Project Progress Report","recovery":"Recovery and Mitigation Report","variation":"Variation Impact Report"}.get(rt,"Integrated Project Controls Report"),detect_data_date(corpus,sch),"",src,sch,comps,ev,metrics,prog,milestones(src,sch),list_items(corpus,["constraint","issue","blocker","hold point"]),list_items(corpus,["risk","threat","exposure"]),list_items(corpus,["action","recommendation","next step"]),conclusions,list(dict.fromkeys(warnings)))
    if context: apply_context(model,context)
    return model

def apply_context(m:Model,c:Dict[str,Any]):
    for k in ["project_name","title","data_date","period"]:
        if c.get(k):setattr(m,k,c[k])
    if c.get("report_type"):
        m.report_type="eot" if c["report_type"]=="tia" else c["report_type"]
    for k in ["metrics","progress"]:
        if isinstance(c.get(k),dict):getattr(m,k).update(c[k])
    for k in ["milestones","constraints","risks","actions","conclusions","warnings"]:
        if isinstance(c.get(k),list):setattr(m,k,c[k])
    if isinstance(c.get("events"),list):
        out=[]
        for i,x in enumerate(c["events"],1):
            if isinstance(x,dict):
                allowed=set(Event.__dataclass_fields__); y={k:v for k,v in x.items() if k in allowed};y.setdefault("event_id",f"EV{i:02d}");y.setdefault("title",f"Event {i}");out.append(Event(**y))
        if out:m.events=out

# --------------------------- Design ----------------------------------------

def wrap(t:str,w:int)->List[str]:
    out=[]
    for p in s(t).splitlines() or [s(t)]: out += textwrap.wrap(p,width=max(8,w),break_long_words=False) or [""]
    return out

def text_layout(e:E)->Tuple[List[str],float,bool]:
    """Return wrapped lines, an automatically fitted font size, and truncation status."""
    if e.kind!="text":return [],e.fs,False
    start=max(6.0,float(e.fs)); minimum=max(9.0,start*.72); fs=start
    while fs>=minimum-.01:
        chars=max(8,int(e.w/max(fs*.53,1))); rendered=wrap(e.text,chars)
        max_lines=max(1,int(e.h/max(fs*1.22,1))) if e.h else len(rendered)
        if len(rendered)<=max_lines:return rendered,fs,False
        fs-=.5
    fs=minimum; chars=max(8,int(e.w/max(fs*.53,1))); rendered=wrap(e.text,chars)
    max_lines=max(1,int(e.h/max(fs*1.22,1))) if e.h else len(rendered)
    truncated=len(rendered)>max_lines
    if truncated:
        rendered=rendered[:max_lines]
        if rendered:rendered[-1]=(rendered[-1][:-3]+"...") if len(rendered[-1])>3 else "..."
    return rendered,fs,truncated


def lines(e:E)->List[str]:
    return text_layout(e)[0]

class Designer:
    def __init__(self,m:Model,b:Dict[str,Any]):self.m=m;self.b=b;self.t=THEME
    def base(self,kicker,title,status="AUTO-GENERATED"):
        z=Slide(title);a=z.add;t=self.t
        a(E("rect",0,0,W,H,fill=t["bg"],stroke=t["bg"]));a(E("rect",0,0,W,98,fill=t["navy"],stroke=t["navy"]));a(E("line",30,18,0,60,stroke=t["gold"],sw=4))
        a(E("text",52,18,1320,24,kicker.upper(),fs=18,color=t["gold"],bold=True,valign="middle"));a(E("text",52,44,1380,44,title,fs=34,color=t["white"],bold=True,valign="middle"))
        a(E("rect",1505,18,380,64,fill="#071F3B",stroke="#8096AE",sw=2,radius=12));a(E("text",1520,27,350,22,status,fs=17,color=t["gold"],bold=True,align="center",valign="middle"));a(E("text",1520,52,350,22,self.m.data_date or self.m.period or "Current evidence set",fs=16,color=t["white"],bold=True,align="center",valign="middle"))
        a(E("rect",0,1016,W,64,fill="#031E39",stroke="#031E39"));a(E("text",34,1032,760,24,self.b["footer_left"],fs=18,color=t["white"],bold=True,valign="middle"));a(E("text",1110,1032,775,24,self.b["footer_right"],fs=18,color=t["gold"],bold=True,align="right",valign="middle"));return z
    def header(self,z,x,y,w,title,color=None):
        c=color or self.t["navy"];z.add(E("rect",x,y,w,46,fill=c,stroke=c,radius=12));z.add(E("text",x+16,y+7,w-32,30,title,fs=22,color=self.t["white"],bold=True,valign="middle"))
    def kpi(self,z,x,y,w,label,value,sub,color):
        t=self.t;z.add(E("rect",x,y,w,88,fill=t["white"],stroke=t["border"],sw=2,radius=13));z.add(E("text",x+15,y+13,w-30,22,label.upper(),fs=15,color=t["muted"],bold=True,valign="middle"));z.add(E("text",x+15,y+36,w-30,34,value,fs=29,color=color,bold=True,align="right",valign="middle"));z.add(E("text",x+15,y+69,w-30,14,sub,fs=12,color=t["ink"],bold=True,align="right",valign="middle"))
    def points(self,z,x,y,w,items,color=None,fs=19,gap=52):
        c=color or self.t["navy2"]
        for i,q in enumerate(items,1):
            yy=y+(i-1)*gap;z.add(E("circle",x,yy,34,34,fill=c,stroke=c));z.add(E("text",x,yy+1,34,31,str(i),fs=16,color=self.t["white"],bold=True,align="center",valign="middle"));z.add(E("text",x+48,yy-2,w-48,gap-4,q,fs=fs,color=self.t["ink"],valign="top"))
    def build(self)->List[Slide]:
        if self.m.report_type=="progress":return [self.progress_dashboard(),self.progress_chart(),self.milestone_slide(),self.control_slide(),self.conclusion()]
        if self.m.report_type in {"eot","delay","variation"}:return [self.eot_dashboard(),self.calculation()]+[self.event_slide(x) for x in (self.m.events[:8] or [Event("EV01","Evidence Gap","No reliable event register or paired schedules were identified.")])]+[self.conclusion()]
        return [self.hybrid(),self.conclusion()]
    def eot_dashboard(self):
        m,t=self.m,self.t;z=self.base("TIME IMPACT / DELAY ANALYSIS",f"EXECUTIVE DASHBOARD — {m.project_name}",m.report_type.upper()+" POSITION")
        gross=int(m.metrics.get("gross_impact_days",0) or 0);ov=m.metrics.get("overlap_days");net=m.metrics.get("net_eot_days");net=net if net is not None else None
        vals=[("GROSS IMPACT",str(gross) if gross else "TBD","calendar days",t["gold"]),("OVERLAP",str(int(ov)) if ov is not None else "TBD","days to validate",t["teal"]),("NET POSITION",str(int(net)) if net is not None else "TBD","subject to validation",t["green"]),("BASE FINISH",s(m.metrics.get("base_finish") or "Not established"),"before impact",t["navy"]),("IMPACTED FINISH",s(m.metrics.get("latest_finish") or "Not established"),"latest measured",t["gold"]),("SOURCE FILES",str(len(m.sources)),"evidence inputs",t["gray"])]
        cw=(1876-14*5)/6
        for i,v in enumerate(vals):self.kpi(z,22+i*(cw+14),114,cw,*v)
        self.header(z,22,220,1160,"EVENT / IMPACT REGISTER — EVIDENCE-BASED STATUS");z.add(E("rect",22,266,1160,472,fill=t["white"],stroke=t["border"],sw=2))
        headers=["EVENT","WHAT THE EVIDENCE SHOWS","SCHEDULE EFFECT","MOVE","TREATMENT"];widths=[120,370,285,105,280];xx=22
        for w,h in zip(widths,headers):z.add(E("rect",xx,266,w,38,fill=t["navy2"],stroke=t["navy2"]));z.add(E("text",xx+5,272,w-10,26,h,fs=14,color=t["gold"],bold=True,align="center",valign="middle"));xx+=w
        colors=[t["red"],t["purple"],t["orange"],t["gray"]]
        for i in range(4):
            y=304+i*108;x=m.events[i] if i<len(m.events) else None
            if i%2:z.add(E("rect",22,y,1160,108,fill="#F8FAFC",stroke="#F8FAFC"))
            if x:
                c=colors[i];z.add(E("circle",43,y+30,54,54,fill=c,stroke=c));z.add(E("text",43,y+41,54,28,x.event_id,fs=14,color=t["white"],bold=True,align="center",valign="middle"));z.add(E("text",145,y+12,340,85,x.description or x.title,fs=17,color=t["ink"]));z.add(E("text",515,y+12,255,85,x.effect,fs=17,color=t["ink"]));z.add(E("text",800,y+22,95,55,str(x.movement_days) if x.movement_days is not None else "TBD",fs=30,color=c,bold=True,align="center",valign="middle"));z.add(E("text",925,y+12,235,85,x.treatment,fs=16,color=c,bold=True,align="center",valign="middle"))
        z.add(E("rect",1200,220,698,518,fill=t["navy"],stroke=t["navy"],radius=16));z.add(E("text",1228,242,642,34,"ANALYTICAL READING",fs=24,color=t["gold"],bold=True,align="center",valign="middle"))
        read=[f"The engine classified the evidence as {m.title.lower()} with {int(m.confidence*100)}% confidence.",f"{len(m.comparisons)} paired before/after schedule comparison(s) were identified.","Measured schedule movement is separated from responsibility and entitlement.","Concurrency is deducted only where the same Project Finish delay period is duplicated.","Native schedule audit and contract verification remain governing controls."]
        for i,q in enumerate(read,1):yy=295+(i-1)*82;z.add(E("circle",1230,yy,34,34,fill=t["gold"],stroke=t["gold"]));z.add(E("text",1230,yy+1,34,31,str(i),fs=16,color=t["navy"],bold=True,align="center",valign="middle"));z.add(E("text",1278,yy-2,575,62,q,fs=17,color=t["white"]))
        self.header(z,22,762,1876,"MANAGEMENT SUMMARY — WHAT WAS ESTABLISHED");self.points(z,40,826,1820,m.conclusions[:4],fs=18,gap=43);return z
    def calculation(self):
        m,t=self.m,self.t;z=self.base("SCHEDULE IMPACT & CONCURRENCY","GROSS-TO-NET TIME POSITION — SIMPLE, TRACEABLE AND NON-DUPLICATIVE","CALCULATION CONTROL")
        pos=[x for x in m.events if x.movement_days is not None and x.movement_days>0];gross=int(m.metrics.get("gross_impact_days",sum(x.movement_days for x in pos)) or 0);ov=m.metrics.get("overlap_days");net=m.metrics.get("net_eot_days");net=net if net is not None else None
        z.add(E("rect",42,125,880,470,fill=t["navy"],stroke=t["navy"],radius=16));z.add(E("text",70,146,824,35,"CONSOLIDATED TIME CALCULATION",fs=25,color=t["gold"],bold=True,align="center",valign="middle"))
        items=[f"{x.event_id} — {x.title}: {x.movement_days} days" for x in pos[:4]] or ["No measured positive event movement identified."]
        for i,q in enumerate(items,1):y=202+(i-1)*58;z.add(E("circle",76,y,38,38,fill=t["gold"],stroke=t["gold"]));z.add(E("text",76,y+2,38,34,str(i),fs=17,color=t["navy"],bold=True,align="center",valign="middle"));z.add(E("rect",130,y-2,750,43,fill="#0D416F",stroke="#6D89A6",sw=1,radius=8));z.add(E("text",145,y+4,720,30,q,fs=18,color=t["white"],bold=True,valign="middle"))
        z.add(E("rect",88,438,746,52,fill="#0D416F",stroke="#6D89A6",sw=1,radius=10));z.add(E("text",105,446,712,34,f"Gross measured movement = {gross if gross else 'TBD'} calendar days",fs=20,color=t["gold"],bold=True,align="center",valign="middle"));z.add(E("rect",88,504,746,52,fill="#0D416F",stroke="#6D89A6",sw=1,radius=10));z.add(E("text",105,512,712,34,f"Verified net time position = {int(net)} days" if net is not None else "Net time remains open until the source calculation and overlap are validated",fs=19,color=t["green"] if net is not None else t["white"],bold=True,align="center",valign="middle"))
        self.header(z,960,125,918,"MANAGEMENT INTERPRETATION");z.add(E("rect",960,171,918,424,fill=t["white"],stroke=t["border"],sw=2));self.points(z,986,205,855,["Measure each before/after schedule pair separately.","Treat the sum of positive movements as gross only after confirming that the schedule comparisons are independent and non-cumulative.","Deduct only a validated duplicated Project Finish delay period.","Do not add or deduct the same overlap more than once.","Final entitlement remains subject to the contract and formal determination."],fs=18,gap=72)
        self.header(z,42,625,1836,"CONCURRENCY / OVERLAP VISUAL CONTROL");z.add(E("rect",42,671,1836,264,fill=t["white"],stroke=t["border"],sw=2));z.add(E("line",150,790,1540,0,stroke=t["navy"],sw=4))
        if len(pos)>=2:
            z.add(E("rect",300,735,690,22,fill=t["red"],stroke=t["red"]));z.add(E("text",300,699,690,30,f"{pos[0].event_id} — {pos[0].movement_days} days",fs=18,color=t["red"],bold=True,align="center",valign="middle"));z.add(E("rect",670,825,720,22,fill=t["purple"],stroke=t["purple"]));z.add(E("text",670,855,720,30,f"{pos[1].event_id} — {pos[1].movement_days} days",fs=18,color=t["purple"],bold=True,align="center",valign="middle"));ow=300 if ov is None else max(130,min(500,int(ov)*4));z.add(E("rect",670,720,ow,150,fill="#9CDCE0",stroke=t["teal"],sw=3,radius=12,opacity=.55));z.add(E("text",670,884,ow,30,f"{int(ov)}-day overlap" if ov is not None else "Overlap to validate",fs=19,color=t["teal"],bold=True,align="center",valign="middle"))
        else:z.add(E("text",220,755,1500,80,"At least two positive impact events are required to visualize concurrency.",fs=23,color=t["muted"],bold=True,align="center",valign="middle"))
        z.add(E("rect",60,952,1800,44,fill=t["pale_green"],stroke=t["green"],sw=2,radius=10));z.add(E("text",80,959,1760,28,f"CURRENT NET TIME POSITION: {int(net)} CALENDAR DAYS" if net is not None else "NET TIME POSITION: REQUIRES VALIDATED OVERLAP ANALYSIS",fs=20,color=t["green"],bold=True,align="center",valign="middle"));return z
    def event_slide(self,x:Event):
        t=self.t;c=t["orange"] if x.movement_days==0 else t["gray"] if x.movement_days is None else t["purple"] if x.event_id.endswith("2") else t["red"];p=t["pale_orange"] if c==t["orange"] else t["pale_gray"] if c==t["gray"] else t["pale_purple"] if c==t["purple"] else t["pale_red"]
        z=self.base("DETAILED EVENT ANALYSIS",f"{x.event_id} — {x.title}",x.treatment.upper()[:28]);cards=[("WHAT HAPPENED",x.description or "Event description not established."),("ANALYTICAL WORK PERFORMED","The engine mapped available narrative, schedule, table, and date evidence; paired schedules where available; and separated measured time effect from contractual responsibility."),("MEASURED / CURRENT EFFECT",x.effect or x.treatment)]
        for i,(h,b) in enumerate(cards):xx=24+i*622;z.add(E("rect",xx,118,604,205,fill=t["white"],stroke=c if i==1 else t["border"],sw=3 if i==1 else 2,radius=16));z.add(E("text",xx+22,138,560,30,h,fs=20,color=c if i else t["navy"],bold=True));z.add(E("text",xx+22,178,560,125,b,fs=20,color=t["ink"]))
        self.header(z,24,345,1872,"CAUSE — EVIDENCE — SCHEDULE EFFECT — CONCLUSION");boxes=[("CAUSE",x.cause),("EVIDENCE","; ".join(x.evidence[:3]) or "Evidence references require confirmation"),("SCHEDULE EFFECT",x.effect),("TREATMENT",x.treatment)]
        for i,(h,b) in enumerate(boxes):xx=48+i*475;z.add(E("rect",xx,415,390,150,fill=p if i==3 else t["white"],stroke=c if i==3 else t["border"],sw=3 if i==3 else 2,radius=15));z.add(E("text",xx+18,432,354,28,h,fs=19,color=c if i==3 else t["navy"],bold=True,align="center",valign="middle"));z.add(E("text",xx+20,470,350,82,b,fs=17,color=t["ink"],align="center",valign="middle"));
        self.header(z,24,600,1180,"SCHEDULE / IMPACT RESULT");z.add(E("rect",24,646,1180,285,fill=t["white"],stroke=t["border"],sw=2));comp=next((q for q in self.m.comparisons if q.movement_days==x.movement_days),None);bf=comp.before_finish if comp else "Not established";af=comp.after_finish if comp else "Not established";vals=[("BEFORE IMPACT",bf,t["navy"],"#EDF4FA"),("AFTER IMPACT",af,c,p),("PROJECT FINISH MOVE",str(x.movement_days) if x.movement_days is not None else "TBD",t["white"],t["navy"])];positions=[(50,690,330),(425,690,330),(800,690,365)]
        for (h,v,co,fi),(xx,yy,ww) in zip(vals,positions):z.add(E("rect",xx,yy,ww,185,fill=fi,stroke=c if h!="BEFORE IMPACT" else t["border"],sw=3 if h!="BEFORE IMPACT" else 2,radius=14));z.add(E("text",xx+15,yy+18,ww-30,28,h,fs=18,color=t["gold"] if fi==t["navy"] else t["navy"],bold=True,align="center",valign="middle"));z.add(E("text",xx+15,yy+65,ww-30,75,v,fs=58 if h=="PROJECT FINISH MOVE" else 30,color=co,bold=True,align="center",valign="middle"));z.add(E("text",xx+15,yy+145,ww-30,24,"CALENDAR DAYS" if h=="PROJECT FINISH MOVE" else "Project Finish",fs=16,color=t["white"] if fi==t["navy"] else t["muted"],bold=True,align="center",valign="middle"))
        self.header(z,1222,600,674,"MANAGEMENT CONCLUSION",c);z.add(E("rect",1222,646,674,285,fill=p,stroke=c,sw=3));q="The current evidence does not establish reliable Project Finish movement. Further activity-level linkage and native schedule audit are required." if x.movement_days is None else "The event was tested without artificially adding time. The current schedule comparison produces zero Project Finish movement; separate cost or disruption rights may still require assessment." if x.movement_days==0 else f"The current before/after schedule evidence indicates {x.movement_days} calendar days of Project Finish movement. Reconcile this with overlap, mitigation, responsibility, notice compliance, and the contract before stating final entitlement.";z.add(E("text",1250,680,618,155,q,fs=21,color=t["ink"]));z.add(E("text",1250,850,618,24,"EVIDENCE POSITION",fs=16,color=t["navy"],bold=True));z.add(E("text",1250,878,618,38,f"Confidence: {x.confidence}. Responsibility and entitlement remain subject to verification.",fs=17,color=c,bold=True));z.add(E("rect",24,950,1872,46,fill=t["pale_green"],stroke=t["green"],sw=2,radius=10));z.add(E("text",42,958,1836,28,f"{x.event_id}: {x.movement_days} calendar days measured movement" if x.movement_days is not None else f"{x.event_id}: movement not established",fs=19,color=t["green"],bold=True,align="center",valign="middle"));return z
    def progress_dashboard(self):
        m,t=self.m,self.t;z=self.base("PROJECT PROGRESS REPORT",f"EXECUTIVE PROGRESS DASHBOARD — {m.project_name}","PERFORMANCE STATUS");pl=m.progress.get("planned");ac=m.progress.get("actual");va=m.progress.get("variance");spi=m.progress.get("spi_proxy");vals=[("PLANNED",f"{pl:.2f}%" if pl is not None else "TBD","latest extracted",t["navy"]),("ACTUAL",f"{ac:.2f}%" if ac is not None else "TBD","latest extracted",t["gold"]),("VARIANCE",f"{va:+.2f}%" if va is not None else "TBD","actual minus planned",t["green"] if va is not None and va>=0 else t["red"]),("SPI",f"{spi:.3f}" if spi is not None else "TBD","schedule performance",t["green"] if spi is not None and spi>=1 else t["red"]),("FORECAST FINISH",s(m.metrics.get("latest_finish") or "Not established"),"latest schedule",t["gold"]),("CONSTRAINTS",str(len(m.constraints)),"extracted items",t["orange"])]
        cw=(1876-14*5)/6
        for i,v in enumerate(vals):self.kpi(z,22+i*(cw+14),114,cw,*v)
        self.header(z,22,220,1180,"PERFORMANCE READING");z.add(E("rect",22,266,1180,470,fill=t["white"],stroke=t["border"],sw=2));self.points(z,52,306,1100,(m.conclusions[:2]+[f"{len(m.milestones)} milestone record(s) identified.",f"{len(m.constraints)} constraint/issue item(s) and {len(m.risks)} risk item(s) extracted.","Forecast dates and critical activities require native schedule verification."])[:5],fs=20,gap=78)
        self.header(z,1230,220,648,"MANAGEMENT FOCUS");z.add(E("rect",1230,266,648,470,fill=t["navy"],stroke=t["navy"]));focus=["Protect critical and near-critical work fronts.","Close information, access, procurement, and interface constraints.","Reconcile physical progress, schedule progress, and quantities.","Link recovery actions to measurable production and schedule outcomes.","Maintain a four-week look-ahead and decision register."]
        for i,q in enumerate(focus,1):yy=302+(i-1)*78;z.add(E("circle",1255,yy,34,34,fill=t["gold"],stroke=t["gold"]));z.add(E("text",1255,yy+1,34,31,str(i),fs=16,color=t["navy"],bold=True,align="center",valign="middle"));z.add(E("text",1305,yy-3,530,55,q,fs=18,color=t["white"]))
        self.header(z,22,762,1856,"EVIDENCE BASIS");z.add(E("rect",22,808,1856,180,fill=t["white"],stroke=t["border"],sw=2));items=[f"{sum(x.kind=='pdf' for x in m.sources)} PDF file(s)",f"{sum(x.kind=='docx' for x in m.sources)} Word file(s)",f"{sum(x.kind=='excel' for x in m.sources)} Excel file(s)",f"{len(m.schedules)} XER schedule(s)",f"{sum(len(x.tables) for x in m.sources)} structured table(s)",f"{len(m.warnings)} validation warning(s)"]
        for i,q in enumerate(items):
            xx=50+(i%3)*610;yy=835+(i//3)*64
            z.add(E("circle",xx,yy,30,30,fill=t["navy2"],stroke=t["navy2"]))
            z.add(E("text",xx,yy+1,30,28,str(i+1),fs=14,color=t["white"],bold=True,align="center",valign="middle"))
            z.add(E("text",xx+42,yy-1,535,34,q,fs=18,color=t["ink"],valign="middle"))
        return z
    def progress_chart(self):
        m,t=self.m,self.t;z=self.base("PROGRESS PERFORMANCE","PLANNED VERSUS ACTUAL — VARIANCE AND MANAGEMENT RESPONSE","PROGRESS CONTROL");pl=float(m.progress.get("planned",0) or 0);ac=float(m.progress.get("actual",0) or 0);va=ac-pl
        self.header(z,42,125,1150,"PLANNED / ACTUAL PERFORMANCE");z.add(E("rect",42,171,1150,500,fill=t["white"],stroke=t["border"],sw=2));z.add(E("line",140,560,920,0,stroke=t["navy"],sw=4));z.add(E("line",140,260,0,300,stroke=t["navy"],sw=4));mx=max(100,pl,ac);ph=300*pl/mx;ah=300*ac/mx;z.add(E("rect",330,560-ph,190,ph,fill=t["navy2"],stroke=t["navy2"],radius=8));z.add(E("rect",680,560-ah,190,ah,fill=t["gold"],stroke=t["gold"],radius=8));z.add(E("text",330,580,190,32,"Planned",fs=20,color=t["navy"],bold=True,align="center"));z.add(E("text",680,580,190,32,"Actual",fs=20,color=t["navy"],bold=True,align="center"));z.add(E("text",330,510-ph,190,38,f"{pl:.2f}%",fs=28,color=t["navy"],bold=True,align="center"));z.add(E("text",680,510-ah,190,38,f"{ac:.2f}%",fs=28,color=t["gold"],bold=True,align="center"));z.add(E("rect",160,602,910,46,fill=t["pale_green"] if va>=0 else t["pale_red"],stroke=t["green"] if va>=0 else t["red"],sw=2,radius=10));z.add(E("text",175,610,880,28,f"PROGRESS VARIANCE = {va:+.2f} PERCENTAGE POINTS",fs=20,color=t["green"] if va>=0 else t["red"],bold=True,align="center",valign="middle"))
        self.header(z,1220,125,658,"CONTROL RESPONSE");z.add(E("rect",1220,171,658,500,fill=t["white"],stroke=t["border"],sw=2));self.points(z,1248,205,590,["Confirm the latest approved baseline and data date.","Reconcile physical quantities with activity status.","Identify critical and near-critical delayed activities.","Assign owners and closure dates to constraints.","Update look-ahead and recovery actions with measurable outputs.","Verify forecast dates in the native schedule."],fs=18,gap=70)
        self.header(z,42,704,1836,"AUTOMATED INTERPRETATION");z.add(E("rect",42,750,1836,236,fill=t["navy"],stroke=t["navy"]));
        for i,q in enumerate(m.conclusions[:4],1):
            yy=782+(i-1)*48
            z.add(E("circle",75,yy,34,34,fill=t["gold"],stroke=t["gold"]))
            z.add(E("text",75,yy+1,34,31,str(i),fs=16,color=t["navy"],bold=True,align="center",valign="middle"))
            z.add(E("text",125,yy-3,1695,44,q,fs=18,color=t["white"],valign="middle"))
        return z
    def milestone_slide(self):
        m,t=self.m,self.t;z=self.base("MILESTONE CONTROL","MILESTONE STATUS — DATES, SOURCE TRACEABILITY AND MANAGEMENT ATTENTION","MILESTONE REVIEW");self.header(z,42,125,1836,"EXTRACTED MILESTONE REGISTER");z.add(E("rect",42,171,1836,650,fill=t["white"],stroke=t["border"],sw=2));headers=["NO.","MILESTONE / ACTIVITY","DATE","STATUS","SOURCE"];widths=[90,760,260,260,466];xx=42
        for w,h in zip(widths,headers):z.add(E("rect",xx,171,w,44,fill=t["navy2"],stroke=t["navy2"]));z.add(E("text",xx+5,178,w-10,30,h,fs=16,color=t["gold"],bold=True,align="center",valign="middle"));xx+=w
        for i in range(8):
            y=215+i*72;row=m.milestones[i] if i<len(m.milestones) else None
            if i%2:z.add(E("rect",42,y,1836,72,fill="#F8FAFC",stroke="#F8FAFC"))
            if row:
                vals=[str(i+1),row.get("milestone",""),row.get("date",""),row.get("status",""),row.get("source","")];xx=42
                for j,(w,v) in enumerate(zip(widths,vals)):z.add(E("text",xx+10,y+8,w-20,56,s(v),fs=17 if j!=1 else 18,color=t["ink"],bold=j in {0,1},align="center" if j in {0,2,3} else "left",valign="middle"));xx+=w
        z.add(E("rect",60,850,1800,130,fill=t["pale_green"],stroke=t["green"],sw=2,radius=12));z.add(E("text",88,870,1744,90,"Milestone dates extracted from schedules and registers must be reconciled to the approved baseline, latest update, contractual requirements, and any formally determined extension before issue.",fs=21,color=t["green"],bold=True,align="center",valign="middle"));return z
    def control_slide(self):
        m,t=self.m,self.t;z=self.base("RISKS, CONSTRAINTS & ACTIONS","MANAGEMENT CONTROL REGISTER — PRIORITIES, OWNERSHIP AND CLOSURE","ACTION CONTROL");groups=[("CONSTRAINTS / ISSUES",m.constraints,t["red"],t["pale_red"]),("RISKS",m.risks,t["orange"],t["pale_orange"]),("ACTIONS / RECOMMENDATIONS",m.actions,t["green"],t["pale_green"])]
        for (h,items,c,p),xx in zip(groups,[24,646,1268]):self.header(z,xx,125,604,h,c);z.add(E("rect",xx,171,604,780,fill=t["white"],stroke=t["border"],sw=2));
        for (h,items,c,p),xx in zip(groups,[24,646,1268]):
            if not items:z.add(E("text",xx+30,220,544,120,f"No structured {h.lower()} were identified.",fs=20,color=t["muted"],align="center",valign="middle"))
            for i,q in enumerate(items[:8],1):yy=205+(i-1)*90;z.add(E("rect",xx+24,yy,556,78,fill=p,stroke=c,sw=2,radius=10));z.add(E("circle",xx+38,yy+22,34,34,fill=c,stroke=c));z.add(E("text",xx+38,yy+23,34,31,str(i),fs=15,color=t["white"],bold=True,align="center",valign="middle"));z.add(E("text",xx+86,yy+9,474,60,q,fs=16,color=t["ink"],valign="middle"))
        z.add(E("rect",42,968,1836,34,fill=t["navy"],stroke=t["navy"],radius=8));z.add(E("text",55,973,1810,24,"EVERY MATERIAL CONSTRAINT AND ACTION SHOULD HAVE AN OWNER, REQUIRED-BY DATE, STATUS, AND VERIFIED SCHEDULE / PRODUCTION EFFECT.",fs=16,color=t["gold"],bold=True,align="center",valign="middle"));return z
    def hybrid(self):
        m,t=self.m,self.t;z=self.base("INTEGRATED PROJECT CONTROLS",f"EXECUTIVE REPORT — {m.project_name}","AUTO-CLASSIFIED");vals=[("REPORT TYPE",m.report_type.upper(),"automatic classification",t["gold"]),("CONFIDENCE",f"{int(m.confidence*100)}%","classification confidence",t["green"]),("SOURCE FILES",str(len(m.sources)),"parsed inputs",t["navy"]),("XER SNAPSHOTS",str(len(m.schedules)),"schedule inputs",t["purple"]),("EVENTS",str(len(m.events)),"identified items",t["orange"]),("WARNINGS",str(len(m.warnings)),"validation items",t["red"])]
        cw=(1876-14*5)/6
        for i,v in enumerate(vals):self.kpi(z,22+i*(cw+14),114,cw,*v)
        self.header(z,22,220,900,"EXTRACTED PROJECT INTELLIGENCE");z.add(E("rect",22,266,900,650,fill=t["white"],stroke=t["border"],sw=2));items=[f"{sum(x.kind=='pdf' for x in m.sources)} PDF file(s)",f"{sum(x.kind=='docx' for x in m.sources)} Word file(s)",f"{sum(x.kind=='excel' for x in m.sources)} Excel file(s)",f"{len(m.schedules)} XER schedule(s)"]+m.conclusions;self.points(z,50,305,835,items[:10],fs=18,gap=58)
        self.header(z,950,220,928,"VALIDATION & NEXT STEPS");z.add(E("rect",950,266,928,650,fill=t["white"],stroke=t["border"],sw=2));items=(m.warnings or ["No major automated warning recorded."])+["Confirm the intended report type where classification is ambiguous.","Use AI context JSON for scanned documents or reviewed conclusions.","Verify all dates, calculations, and conclusions against native source records."];self.points(z,978,305,860,items[:10],color=t["gold"],fs=18,gap=58);return z
    def conclusion(self):
        m,t=self.m,self.t;z=self.base("MANAGEMENT CONCLUSION & CONTROL","REPORT CONCLUSION — MEASURED FINDINGS, LIMITATIONS AND NEXT DECISIONS","FINAL REVIEW");self.header(z,42,125,1160,"CONSOLIDATED FINDINGS");z.add(E("rect",42,171,1160,485,fill=t["white"],stroke=t["border"],sw=2));self.points(z,70,205,1090,m.conclusions[:7],fs=19,gap=63);self.header(z,1230,125,648,"EVIDENCE & LIMITATIONS");z.add(E("rect",1230,171,648,485,fill=t["white"],stroke=t["border"],sw=2));e=[f"{len(m.sources)} source file(s) parsed",f"{len(m.schedules)} XER schedule(s) parsed",f"{sum(len(x.tables) for x in m.sources)} structured table(s) extracted"]+m.warnings;self.points(z,1252,202,590,e[:8],color=t["gold"],fs=17,gap=54);self.header(z,42,690,1836,"RECOMMENDED DECISION PATH");z.add(E("rect",42,736,1836,218,fill=t["navy"],stroke=t["navy"]));path=["Verify dates, activity logic, calendars, constraints, and longest-path behavior in the native scheduling application.","Confirm event chronology and responsibility against contemporaneous records and the governing contract.","Validate concurrency by activity-level criticality and common Project Finish delay period.","Issue the management or contractual position using the verified result, evidence references, and determination route."]
        for i,q in enumerate(path,1):
            yy=770+(i-1)*45
            z.add(E("circle",75,yy,34,34,fill=t["gold"],stroke=t["gold"]))
            z.add(E("text",75,yy+1,34,31,str(i),fs=16,color=t["navy"],bold=True,align="center",valign="middle"))
            z.add(E("text",125,yy-3,1680,40,q,fs=18,color=t["white"],valign="middle"))
        return z

# --------------------------- Renderers -------------------------------------

def render_svg(z:Slide,p:Path):
    out=[f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{H}" viewBox="0 0 {W} {H}">']
    for e in z.elements:
        op=f' opacity="{e.opacity}"' if e.opacity<1 else ""
        if e.kind=="rect":out.append(f'<rect x="{e.x}" y="{e.y}" width="{e.w}" height="{e.h}" rx="{e.radius}" fill="{e.fill}" stroke="{e.stroke}" stroke-width="{e.sw}"{op}/>')
        elif e.kind=="circle":out.append(f'<ellipse cx="{e.x+e.w/2}" cy="{e.y+e.h/2}" rx="{e.w/2}" ry="{e.h/2}" fill="{e.fill}" stroke="{e.stroke}" stroke-width="{e.sw}"{op}/>')
        elif e.kind=="line":out.append(f'<line x1="{e.x}" y1="{e.y}" x2="{e.x+e.w}" y2="{e.y+e.h}" stroke="{e.stroke}" stroke-width="{e.sw}"{op}/>')
        elif e.kind=="text":
            ll,fs,_=text_layout(e)
            if not ll:continue
            lh=fs*1.22;total=len(ll)*lh;sy=e.y+(e.h-total)/2+fs if e.valign=="middle" else e.y+fs;x=e.x+e.w/2 if e.align=="center" else e.x+e.w if e.align=="right" else e.x;anchor="middle" if e.align=="center" else "end" if e.align=="right" else "start";out.append(f'<text x="{x}" y="{sy}" text-anchor="{anchor}" style="font-family:Arial,sans-serif;font-size:{fs}px;font-weight:{700 if e.bold else 400};fill:{e.color};"{op}>')
            for i,q in enumerate(ll):out.append(f'<tspan x="{x}" dy="{0 if i==0 else lh}">{xml_escape(q)}</tspan>')
            out.append('</text>')
    out.append('</svg>');p.parent.mkdir(parents=True,exist_ok=True);p.write_text(''.join(out),encoding='utf-8')

def rgb(c):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(c.lstrip('#').upper())
def ix(x):
    from pptx.util import Inches
    return Inches(x/W*13.333333)
def iy(y):
    from pptx.util import Inches
    return Inches(y/H*7.5)

def render_ppt(slides:List[Slide],p:Path):
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN,MSO_ANCHOR
        from pptx.util import Inches,Pt
    except ImportError:raise RuntimeError('Install python-pptx')
    prs=Presentation();prs.slide_width=Inches(13.333333);prs.slide_height=Inches(7.5)
    for z in slides:
        sl=prs.slides.add_slide(prs.slide_layouts[6])
        for e in z.elements:
            if e.kind in {'rect','circle'}:
                sh=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if e.kind=='rect' and e.radius else MSO_SHAPE.RECTANGLE if e.kind=='rect' else MSO_SHAPE.OVAL,ix(e.x),iy(e.y),ix(e.w),iy(e.h));sh.fill.solid();sh.fill.fore_color.rgb=rgb(e.fill);sh.fill.transparency=int((1-e.opacity)*100)
                if e.sw:sh.line.color.rgb=rgb(e.stroke);sh.line.width=Pt(max(.5,e.sw*.75))
                else:sh.line.fill.background()
            elif e.kind=='line':
                sh=sl.shapes.add_connector(1,ix(e.x),iy(e.y),ix(e.x+e.w),iy(e.y+e.h));sh.line.color.rgb=rgb(e.stroke);sh.line.width=Pt(max(.5,e.sw*.75))
            elif e.kind=='text':
                ll,fs,_=text_layout(e)
                sh=sl.shapes.add_textbox(ix(e.x),iy(e.y),ix(e.w),iy(e.h));tf=sh.text_frame;tf.clear();tf.word_wrap=True;tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=Pt(0);tf.vertical_anchor={'top':MSO_ANCHOR.TOP,'middle':MSO_ANCHOR.MIDDLE,'bottom':MSO_ANCHOR.BOTTOM}.get(e.valign,MSO_ANCHOR.TOP)
                for i,q in enumerate(ll or [e.text]):
                    pa=tf.paragraphs[0] if i==0 else tf.add_paragraph();pa.text=q;pa.alignment={'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT}.get(e.align,PP_ALIGN.LEFT);pa.space_after=pa.space_before=Pt(0);pa.line_spacing=1
                    for r in pa.runs:r.font.name='Arial';r.font.size=Pt(max(6,fs*.75));r.font.bold=e.bold;r.font.color.rgb=rgb(e.color)
    p.parent.mkdir(parents=True,exist_ok=True);prs.save(p)

def render_pngs(svg:List[Path],out:Path)->List[Path]:
    try:import cairosvg
    except ImportError:raise RuntimeError('Install cairosvg')
    out.mkdir(parents=True,exist_ok=True);r=[]
    for x in svg:y=out/(x.stem+'.png');cairosvg.svg2png(url=str(x),write_to=str(y),output_width=W,output_height=H);r.append(y)
    return r

def render_png_ppt(png:List[Path],p:Path):
    from pptx import Presentation
    from pptx.util import Inches
    prs=Presentation();prs.slide_width=Inches(13.333333);prs.slide_height=Inches(7.5)
    for x in png:sl=prs.slides.add_slide(prs.slide_layouts[6]);sl.shapes.add_picture(str(x),0,0,width=prs.slide_width,height=prs.slide_height)
    p.parent.mkdir(parents=True,exist_ok=True);prs.save(p)

def render_pdf(png:List[Path],p:Path):
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4,landscape
    except ImportError:raise RuntimeError('Install reportlab')
    pw,ph=landscape(A4);p.parent.mkdir(parents=True,exist_ok=True);c=canvas.Canvas(str(p),pagesize=(pw,ph),pageCompression=1)
    for x in png:c.drawImage(str(x),0,0,width=pw,height=ph,preserveAspectRatio=False,mask='auto');c.showPage()
    c.save()

def create_contact_sheet(png:List[Path],p:Path,columns:int=2):
    from PIL import Image
    if not png:return
    thumb_w=960;thumb_h=540;rows=(len(png)+columns-1)//columns
    sheet=Image.new("RGB",(thumb_w*columns,thumb_h*rows),(245,248,252))
    for i,path in enumerate(png):
        im=Image.open(path).convert("RGB").resize((thumb_w,thumb_h))
        sheet.paste(im,((i%columns)*thumb_w,(i//columns)*thumb_h))
    p.parent.mkdir(parents=True,exist_ok=True);sheet.save(p,quality=95)


def create_html_gallery(png:List[Path],p:Path):
    tabs=[];panels=[]
    for i,path in enumerate(png):
        label=path.stem.replace("_"," ");active=" active" if i==0 else ""
        tabs.append(f'<button class="tab{active}" onclick="showTab({i})">{xml_escape(label)}</button>')
        panels.append(f'<section id="panel-{i}" class="panel{active}"><img src="../PNG_HIGH_RES/{path.name}" alt="{xml_escape(label)}"></section>')
    template="""<!doctype html><html><head><meta charset="utf-8"><title>Universal Project Report</title><style>
*{box-sizing:border-box}html,body{margin:0;background:#eaf0f6;font-family:Arial,sans-serif;color:#102a43}
header{background:#06294F;color:white;padding:18px 24px}h1{margin:0;font-size:24px}
.tabs{display:flex;gap:8px;padding:12px;background:white;position:sticky;top:0;z-index:5;overflow-x:auto;box-shadow:0 2px 10px #0002}
.tab{border:1px solid #bed0e2;background:#f6f9fc;color:#06294F;padding:11px 14px;border-radius:9px;font-weight:700;cursor:pointer;white-space:nowrap}
.tab.active{background:#06294F;color:#F8B915}main{padding:16px;max-width:1920px;margin:auto}
.panel{display:none;background:white;padding:8px;border-radius:14px;box-shadow:0 4px 18px #0002}.panel.active{display:block}img{display:block;width:100%;border-radius:8px}
</style><script>function showTab(i){document.querySelectorAll('.tab').forEach((e,j)=>e.classList.toggle('active',i===j));document.querySelectorAll('.panel').forEach((e,j)=>e.classList.toggle('active',i===j));}</script></head><body><header><h1>Universal Project Report</h1></header><div class="tabs">__TABS__</div><main>__PANELS__</main></body></html>"""
    html=template.replace("__TABS__","".join(tabs)).replace("__PANELS__","".join(panels))
    p.parent.mkdir(parents=True,exist_ok=True);p.write_text(html,encoding="utf-8")


def _intersect(a:E,b:E)->Tuple[float,float,float,float,float]:
    x1=max(a.x,b.x);y1=max(a.y,b.y);x2=min(a.x+a.w,b.x+b.w);y2=min(a.y+a.h,b.y+b.h)
    if x2<=x1 or y2<=y1:return x1,y1,x2,y2,0.0
    return x1,y1,x2,y2,(x2-x1)*(y2-y1)


def _contains(a:E,b:E,tol:float=1.0)->bool:
    return a.x<=b.x+tol and a.y<=b.y+tol and a.x+a.w>=b.x+b.w-tol and a.y+a.h>=b.y+b.h-tol


def validate(slides:List[Slide])->Dict[str,Any]:
    out={"status":"PASS","slides":[],"errors":0,"warnings":0,"text_overlaps":0,"shape_overlaps":0,"truncated_text":0}
    for i,z in enumerate(slides,1):
        err=[];warn=[];texts=[];shapes=[]
        for j,e in enumerate(z.elements,1):
            if e.x<-.5 or e.y<-.5 or e.x+e.w>W+.5 or e.y+e.h>H+.5:err.append(f"Element {j} outside slide")
            if e.kind=="text":
                ll,fs,truncated=text_layout(e);texts.append((j,e))
                if truncated:warn.append(f"Text element {j} was truncated after automatic font fitting");out["truncated_text"]+=1
                if not ll and e.text.strip():warn.append(f"Text element {j} produced no rendered lines")
            elif e.kind in {"rect","circle"} and e.opacity>=.9 and e.w>0 and e.h>0:shapes.append((j,e))
        for a in range(len(texts)):
            ja,ea=texts[a]
            for jb,eb in texts[a+1:]:
                x1,y1,x2,y2,area=_intersect(ea,eb);iw=max(0,x2-x1);ih=max(0,y2-y1)
                if iw>4 and ih>4 and area>max(80,.08*min(ea.w*ea.h,eb.w*eb.h)):
                    err.append(f"Text elements {ja} and {jb} overlap");out["text_overlaps"]+=1
        for a in range(len(shapes)):
            ja,ea=shapes[a]
            for jb,eb in shapes[a+1:]:
                if _contains(ea,eb) or _contains(eb,ea):continue
                *_,area=_intersect(ea,eb)
                if area>.22*min(ea.w*ea.h,eb.w*eb.h):
                    warn.append(f"Shape elements {ja} and {jb} materially overlap");out["shape_overlaps"]+=1
        row={"slide":i,"title":z.title,"errors":list(dict.fromkeys(err)),"warnings":list(dict.fromkeys(warn))}
        out["slides"].append(row);out["errors"]+=len(row["errors"]);out["warnings"]+=len(row["warnings"])
    if out["errors"]:out["status"]="FAIL"
    elif out["warnings"]:out["status"]="PASS_WITH_WARNINGS"
    return out

# --------------------------- Public API ------------------------------------

def analyze_inputs(input_files:Sequence[str|Path],report_type:str='auto',context:Optional[Dict[str,Any]]=None,strict:bool=False)->Model:
    with tempfile.TemporaryDirectory() as td:
        files=collect_inputs(input_files,Path(td));src=[];sch=[]
        for i,p in enumerate(files,1):
            try:
                a,b=read_source(p);a.meta.setdefault("source_id",f"SRC-{i:03d}");src.append(a);sch += [b] if b else []
            except Exception as e:src.append(Source(str(p),p.suffix[1:],p.stem,meta={"source_id":f"SRC-{i:03d}"},warnings=[f'Reader failed: {e}']))
        if strict and any(any("Reader failed:" in w for w in x.warnings) for x in src):raise RuntimeError("One or more evidence readers failed in strict mode.")
        return build_model(src,sch,report_type,context)

def generate_report(input_files:Sequence[str|Path],output_directory:str|Path,report_type:str='auto',config_path:Optional[str|Path]=None,context:Optional[Dict[str,Any]]=None,context_path:Optional[str|Path]=None,strict:bool=False,keep_working:bool=False)->Dict[str,Any]:
    if context_path:context=load_json(Path(context_path))
    out=Path(output_directory).expanduser().resolve();out.mkdir(parents=True,exist_ok=True);work=out/'_working';shutil.rmtree(work,ignore_errors=True);work.mkdir()
    files=collect_inputs(input_files,work);src=[];sch=[]
    for i,p in enumerate(files,1):
        try:a,b=read_source(p);a.meta.setdefault("source_id",f"SRC-{i:03d}");src.append(a);sch += [b] if b else []
        except Exception as e:src.append(Source(str(p),p.suffix[1:],p.stem,meta={"source_id":f"SRC-{i:03d}"},warnings=[f'Reader failed: {e}']))
    if strict and any(any("Reader failed:" in w for w in x.warnings) for x in src):raise RuntimeError("One or more evidence readers failed in strict mode.")
    m=build_model(src,sch,report_type,context);slides=Designer(m,branding(config_path)).build();val=validate(slides)
    if strict and val["status"]=="FAIL":raise RuntimeError("Layout validation failed in strict mode: "+json.dumps(val,ensure_ascii=False))
    svgdir=out/'SVG_EDITABLE';pngdir=out/'PNG_HIGH_RES';pptdir=out/'POWERPOINT';pdfdir=out/'PDF_A4_LANDSCAPE';datadir=out/'PROJECT_INTELLIGENCE';vdir=out/'VALIDATION';previewdir=out/'PREVIEW';htmldir=out/'HTML'
    for x in [svgdir,pngdir,pptdir,pdfdir,datadir,vdir,previewdir,htmldir]:x.mkdir(parents=True,exist_ok=True)
    svgs=[]
    for i,z in enumerate(slides,1):p=svgdir/f'{i:02d}_{slug(z.title)}.svg';render_svg(z,p);svgs.append(p)
    editable=pptdir/'Universal_Project_Report_EDITABLE.pptx';render_ppt(slides,editable);pngs=render_pngs(svgs,pngdir);pngppt=pptdir/'Universal_Project_Report_PNG.pptx';render_png_ppt(pngs,pngppt);pdf=pdfdir/'Universal_Project_Report_A4_Landscape_FULL_BLEED.pdf';render_pdf(pngs,pdf)
    contact=previewdir/'Universal_Project_Report_Contact_Sheet.png';create_contact_sheet(pngs,contact);gallery=htmldir/'Universal_Project_Report_Gallery.html';create_html_gallery(pngs,gallery)
    modelp=datadir/'project_model.json';save_json(modelp,asdict(m));inventoryp=datadir/'source_inventory.json';save_json(inventoryp,source_inventory(m.sources));templatep=datadir/'AI_CONTEXT_TEMPLATE.json';save_json(templatep,context_template())
    valp=vdir/'layout_and_content_validation.json';save_json(valp,val);(vdir/'REPORT_GENERATION_SUMMARY.txt').write_text(f'Engine version: {VERSION}\nPython: {platform.python_version()}\nProject: {m.project_name}\nReport type: {m.report_type}\nClassification confidence: {m.confidence}\nSources: {len(m.sources)}\nSlides: {len(slides)}\nValidation: {val["status"]}\nLayout errors: {val["errors"]}\nLayout warnings: {val["warnings"]}\nContent warnings: {len(m.warnings)}\n',encoding='utf-8')
    if not keep_working:shutil.rmtree(work,ignore_errors=True)
    manifestp=out/'ENGINE_RUN_MANIFEST.json';manifest={"engine_version":VERSION,"generated_at":datetime.now().isoformat(timespec="seconds"),"project_name":m.project_name,"report_type":m.report_type,"classification_confidence":m.confidence,"validation_status":val["status"],"files":[]}
    for f in sorted(out.rglob('*')):
        if f.is_file() and f!=manifestp:manifest["files"].append({"path":str(f.relative_to(out)),"size_bytes":f.stat().st_size,"sha256":file_sha256(f)})
    save_json(manifestp,manifest)
    zp=out.parent/f'{out.name}_FULL_PACKAGE.zip';zp.unlink(missing_ok=True)
    with zipfile.ZipFile(zp,'w',zipfile.ZIP_DEFLATED) as z:
        for f in out.rglob('*'):
            if f.is_file():z.write(f,f.relative_to(out.parent))
    return {'status':'completed','engine_version':VERSION,'project_name':m.project_name,'report_type':m.report_type,'confidence':m.confidence,'editable_powerpoint':str(editable),'png_powerpoint':str(pngppt),'pdf':str(pdf),'svg_directory':str(svgdir),'png_directory':str(pngdir),'html_gallery':str(gallery),'contact_sheet':str(contact),'project_model':str(modelp),'source_inventory':str(inventoryp),'context_template':str(templatep),'validation':str(valp),'manifest':str(manifestp),'package_zip':str(zp),'slide_count':len(slides),'warnings':m.warnings}

# --------------------------- CLI -------------------------------------------

def main(argv=None):
    ap=argparse.ArgumentParser(description=f"Universal project report engine v{VERSION}")
    ap.add_argument("--version",action="version",version=VERSION)
    sp=ap.add_subparsers(dest="cmd")
    g=sp.add_parser("generate",help="Analyze evidence and generate the complete report package")
    g.add_argument("--input",nargs="+",required=True);g.add_argument("--output",required=True);g.add_argument("--report-type",default="auto",choices=sorted(ALLOWED_REPORT_TYPES));g.add_argument("--config");g.add_argument("--context-json");g.add_argument("--strict",action="store_true");g.add_argument("--keep-working",action="store_true")
    a=sp.add_parser("analyze",help="Analyze evidence and write only the normalized project model")
    a.add_argument("--input",nargs="+",required=True);a.add_argument("--output",required=True);a.add_argument("--report-type",default="auto",choices=sorted(ALLOWED_REPORT_TYPES));a.add_argument("--context-json");a.add_argument("--strict",action="store_true")
    c=sp.add_parser("context-template",help="Write a reviewed-context JSON template for AI vision/OCR enrichment");c.add_argument("--output",required=True)
    args=ap.parse_args(argv)
    if args.cmd=="generate":
        print(json.dumps(generate_report(args.input,args.output,args.report_type,args.config,context_path=args.context_json,strict=args.strict,keep_working=args.keep_working),indent=2,ensure_ascii=False));return 0
    if args.cmd=="analyze":
        m=analyze_inputs(args.input,args.report_type,load_json(Path(args.context_json)) if args.context_json else None,strict=args.strict);save_json(Path(args.output),asdict(m));print(args.output);return 0
    if args.cmd=="context-template":save_json(Path(args.output),context_template());print(args.output);return 0
    ap.print_help();return 1

if __name__=='__main__':raise SystemExit(main())
