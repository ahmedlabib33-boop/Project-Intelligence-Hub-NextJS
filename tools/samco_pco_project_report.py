"""Data-driven adapter for the supplied SAMCO-PCO 31-slide report engine.

The supplied engine establishes the 31-slide structure and SAMCO visual
language.  Its shipped demo text is deliberately not emitted here: every
published value, table, count, and source hash is read from the active
project's CSV inputs at generation time.
"""

from __future__ import annotations

import hashlib
import html
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import samco_pco_report_generator as samco_template

from project_input_contracts import load_logical_rows, load_payment_rows, logical_source_path, read_csv_rows


WORKSPACE_ROOT = Path(__file__).resolve().parents[1]
ENGINE_SOURCE = WORKSPACE_ROOT / "tools" / "samco_pco_report_generator.py"
REPORT_KEY = "samco_pco_31_slide_report"
REPORT_STEM = "05_samco_pco_31_slide_report"
REPORT_GENERATOR_VERSION = "2026.08.samco-pco-data-driven.v1"

REPORT_SECTIONS: tuple[tuple[str, str, tuple[str, ...]], ...] = (
    ("PRO-TIA-01", "Time Impact & Extension of Time", ("delay_classification", "tia_recovery", "p6")),
    ("PRO-DLY-02", "Delay Analysis", ("delay_events", "delay_classification")),
    ("PRO-PRG-03", "Project Progress", ("activities", "progress")),
    ("PRO-REC-04", "Recovery & Mitigation", ("tia_recovery", "activities")),
    ("PRO-VAR-05", "Variation Report", ("claims",)),
    ("PRO-IPC-06", "Interim Payment", ("payments",)),
    ("PRO-RES-07", "Resource & Manpower", ("activities",)),
    ("PRO-EQP-08", "Equipment & Plant", tuple()),
    ("PRO-CST-09", "Cost Control", ("evm",)),
    ("PRO-EVM-10", "Earned Value", ("evm",)),
    ("PRO-CFS-11", "Cash Flow", ("planned_cash_flow", "payments")),
    ("PRO-SCV-12", "S-Curve", ("s_curve",)),
    ("PRO-PRC-13", "Procurement Status", ("activities",)),
    ("PRO-MAT-14", "Material Status", ("activities",)),
    ("PRO-LAH-15", "Labour Hours", ("activities",)),
    ("PRO-CPM-16", "Critical Path", ("p6", "activities")),
    ("PRO-FLT-17", "Float Analysis", ("p6", "activities")),
    ("PRO-MIL-18", "Milestone Status", ("milestones",)),
    ("PRO-RSK-19", "Risk Report", ("risks",)),
    ("PRO-RFI-20", "RFI & Submittals", ("rfi",)),
    ("PRO-QAQC-21", "QA/QC Performance", tuple()),
    ("PRO-PROD-22", "Productivity", ("activities", "progress")),
    ("PRO-EXE-23", "Executive Report", ("projects", "activities", "evm", "risks")),
    ("PRO-BLC-24", "Baseline vs Current", ("p6", "activities")),
    ("PRO-FCF-25", "Forecast Completion", ("activities", "evm")),
    ("PRO-ML-26", "ML Project Controls", ("activities", "evm", "risks")),
    ("PRO-CAD-27", "Contract Administration", ("contracts", "claims")),
    ("PRO-CHG-28", "Change Control", ("claims",)),
    ("PRO-DOC-29", "Document Control", tuple()),
    ("PRO-IFC-30", "Interface Management", ("ifc",)),
)


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _project_root(project: dict[str, Any]) -> Path:
    return WORKSPACE_ROOT / "projects" / str(project.get("sector") or "") / str(project.get("project_folder_name") or "")


def _source_catalog(project: dict[str, Any]) -> dict[str, tuple[Path, list[dict[str, str]]]]:
    project_root = _project_root(project)
    data_dir = project_root / "01-data" / "import_templates"
    delay_dir = project_root / "02-delay_analysis" / "unified_tia_csv"
    payment_path = data_dir / "payments.csv" if (data_dir / "payments.csv").exists() else delay_dir / "08- payments.csv"
    ordinary = {
        "projects": data_dir / "projects.csv",
        "contracts": data_dir / "contracts.csv",
        "payments": payment_path,
        "planned_cash_flow": data_dir / "planned_cash_flow.csv",
        "risks": data_dir / "risks.csv",
        "claims": data_dir / "claims.csv",
        "milestones": data_dir / "milestones.csv",
        "delay_events": data_dir / "delay_events.csv",
        "s_curve": data_dir / "s_curve.csv",
        "delay_classification": delay_dir / "12_delay_event_classification.csv",
        "tia_recovery": delay_dir / "13_tia_recovery_scenario.csv",
        "rfi": delay_dir / "09- rfi_status.csv",
        "ifc": delay_dir / "07- ifc_conflict.csv",
        "p6": delay_dir / "04- p6_activity_export.csv",
    }
    catalog = {name: (path, read_csv_rows(path)) for name, path in ordinary.items()}
    catalog["payments"] = (payment_path, load_payment_rows(data_dir, delay_dir))
    for logical_name in ("activities", "progress", "evm"):
        canonical_name = "progress_updates" if logical_name == "progress" else logical_name
        catalog[logical_name] = (
            logical_source_path(data_dir, delay_dir, canonical_name),
            load_logical_rows(data_dir, delay_dir, canonical_name),
        )
    return catalog


def _source_label(project_root: Path, path: Path) -> str:
    try:
        return path.relative_to(project_root).as_posix()
    except ValueError:
        return path.name


def _section_data(project: dict[str, Any]) -> list[dict[str, Any]]:
    project_root = _project_root(project)
    catalog = _source_catalog(project)
    sections: list[dict[str, Any]] = []
    for code, title, source_names in REPORT_SECTIONS:
        sources = []
        primary_rows: list[dict[str, str]] = []
        for source_name in source_names:
            path, rows = catalog[source_name]
            if rows and not primary_rows:
                primary_rows = rows
            sources.append({
                "logical_table": source_name,
                "path": _source_label(project_root, path),
                "row_count": len(rows),
                "sha256": _sha256(path) if path.exists() else None,
            })
        sections.append({
            "code": code,
            "title": title,
            "status": "DATA AVAILABLE" if primary_rows else "AWAITING CONTROLLED CSV DATA",
            "sources": sources,
            "rows": primary_rows[:50],
            "total_rows": len(primary_rows),
        })
    return sections


def _table_html(rows: list[dict[str, str]]) -> str:
    if not rows:
        return "<p class='empty'>No controlled CSV rows are available for this report section.</p>"
    headers = list(rows[0].keys())
    head = "".join(f"<th>{html.escape(header)}</th>" for header in headers)
    body = "".join(
        "<tr>" + "".join(f"<td>{html.escape(str(row.get(header, '') or ''))}</td>" for header in headers) + "</tr>"
        for row in rows
    )
    return f"<div class='table-wrap'><table><thead><tr>{head}</tr></thead><tbody>{body}</tbody></table></div>"


def _write_html(path: Path, project: dict[str, Any], sections: list[dict[str, Any]], generated_at: str) -> None:
    tabs = "".join(
        f"<button role='tab' aria-selected={'true' if index == 0 else 'false'} data-tab='{index}'>{html.escape(item['code'])}</button>"
        for index, item in enumerate(sections)
    )
    panels = []
    for index, item in enumerate(sections):
        source_rows = "".join(
            "<li><b>{logical}</b> — {path} — {count} rows — SHA-256 {hash}</li>".format(
                logical=html.escape(source["logical_table"]),
                path=html.escape(source["path"]),
                count=source["row_count"],
                hash=html.escape(source["sha256"] or "not found"),
            )
            for source in item["sources"]
        ) or "<li>No source contract is mapped to this template section.</li>"
        panels.append(
            f"<section class='panel' data-panel='{index}' {'hidden' if index else ''}>"
            f"<header><span>{html.escape(item['code'])}</span><h2>{html.escape(item['title'])}</h2><strong>{html.escape(item['status'])}</strong></header>"
            f"<p>Primary table: {item['total_rows']} controlled rows. Values below are from the selected project only.</p>"
            f"<h3>CSV provenance</h3><ul>{source_rows}</ul><h3>Controlled data preview</h3>{_table_html(item['rows'])}</section>"
        )
    title = f"SAMCO-PCO 31-Slide Project Controls Report — {project.get('project_display_name') or project.get('project_id')}"
    metrics = [
        ("Project ID", project.get("project_id")),
        ("Status", project.get("status")),
        ("Contract value", project.get("contract_value")),
        ("Actual progress", project.get("actual_progress")),
        ("SPI / CPI", f"{project.get('spi')} / {project.get('cpi')}"),
        ("Delay exposure (days)", project.get("delay_days")),
    ]
    metric_html = "".join(f"<div><small>{html.escape(label)}</small><b>{html.escape(str(value if value is not None else 'N/A'))}</b></div>" for label, value in metrics)
    path.write_text(
        "<!doctype html><html><head><meta charset='utf-8'><meta name='viewport' content='width=device-width,initial-scale=1'>"
        f"<title>{html.escape(title)}</title><style>"
        "body{margin:0;background:#f4f7fb;color:#102a43;font:15px/1.5 Arial,sans-serif}.hero{padding:28px 5vw;background:#003366;color:#fff}.hero h1{margin:0 0 6px;font-size:28px}.hero p{margin:0;color:#cfe0f0}.metrics{display:grid;grid-template-columns:repeat(auto-fit,minmax(150px,1fr));gap:10px;margin-top:20px}.metrics div{background:#0d4775;padding:10px;border-radius:8px}.metrics small{display:block;color:#b9d6ec}.metrics b{font-size:17px}.tabs{display:flex;gap:6px;overflow:auto;padding:16px 5vw;background:#e5eef6}.tabs button{white-space:nowrap;border:1px solid #aebfce;background:#fff;padding:8px 10px;border-radius:6px;color:#073b66;font-weight:bold;cursor:pointer}.tabs button[aria-selected=true]{background:#073b66;color:#fff}.panel{padding:24px 5vw}.panel header{display:flex;gap:12px;align-items:baseline;flex-wrap:wrap}.panel h2{margin:0}.panel header span{font-weight:bold;color:#0b69a3}.panel header strong{margin-left:auto;color:#0a7d45}.table-wrap{overflow:auto;background:#fff;border:1px solid #d6e0e8;border-radius:8px}table{border-collapse:collapse;width:100%;min-width:720px}th,td{padding:8px 10px;border-bottom:1px solid #e5edf3;text-align:left;vertical-align:top}th{position:sticky;top:0;background:#e8f0f7}.empty{padding:16px;background:#fff7e6;border-left:4px solid #c78a10}footer{padding:20px 5vw;color:#526f87}</style></head><body>"
        f"<div class='hero'><h1>{html.escape(title)}</h1><p>Data-driven adapter based on the supplied SAMCO-PCO report engine. Generated {html.escape(generated_at)}.</p><div class='metrics'>{metric_html}</div></div>"
        f"<nav class='tabs' role='tablist'>{tabs}</nav>{''.join(panels)}"
        "<footer>Only controlled project CSV values are shown. A schedule/EOT conclusion requires verified Primavera P6 recalculation and entitlement review.</footer>"
        "<script>const tabs=[...document.querySelectorAll('[data-tab]')],panels=[...document.querySelectorAll('[data-panel]')];tabs.forEach(t=>t.onclick=()=>{tabs.forEach(x=>x.setAttribute('aria-selected',String(x===t)));panels.forEach(p=>p.hidden=p.dataset.panel!==t.dataset.tab)})</script>"
        "</body></html>",
        encoding="utf-8",
    )


def _write_pdf(path: Path, project: dict[str, Any], sections: list[dict[str, Any]], generated_at: str) -> None:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    styles = getSampleStyleSheet()
    story: list[Any] = []
    for index, item in enumerate(sections):
        story.append(Paragraph(f"{item['code']} | {html.escape(item['title'])}", styles["Title"]))
        story.append(Paragraph(html.escape(str(project.get("project_display_name") or project.get("project_id") or "")), styles["Heading2"]))
        story.append(Paragraph(f"Status: {html.escape(item['status'])}. Generated {html.escape(generated_at)}.", styles["BodyText"]))
        source_data = [["Logical table", "Project-relative CSV", "Rows", "SHA-256"]]
        source_data.extend([[source["logical_table"], source["path"], str(source["row_count"]), source["sha256"] or "not found"] for source in item["sources"]])
        source_table = Table(source_data, colWidths=[42 * mm, 70 * mm, 18 * mm, 115 * mm])
        source_table.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#003366")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white), ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#aebfce")), ("FONTSIZE", (0, 0), (-1, -1), 7), ("VALIGN", (0, 0), (-1, -1), "TOP")]))
        story.extend([Spacer(1, 6 * mm), source_table, Spacer(1, 6 * mm)])
        if item["rows"]:
            headers = list(item["rows"][0].keys())[:6]
            data = [headers] + [[str(row.get(header, ""))[:55] for header in headers] for row in item["rows"][:12]]
            table = Table(data, repeatRows=1)
            table.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#dce9f4")), ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#cbd5df")), ("FONTSIZE", (0, 0), (-1, -1), 6), ("VALIGN", (0, 0), (-1, -1), "TOP")]))
            story.append(table)
        else:
            story.append(Paragraph("No controlled CSV rows are available for this section.", styles["BodyText"]))
        if index < len(sections) - 1:
            story.append(PageBreak())
    SimpleDocTemplate(str(path), pagesize=landscape(letter), leftMargin=12 * mm, rightMargin=12 * mm, topMargin=12 * mm, bottomMargin=12 * mm).build(story)


def _write_pptx(path: Path, project: dict[str, Any], sections: list[dict[str, Any]], generated_at: str) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    navy = samco_template.COLOR_TITLE
    white = samco_template.COLOR_WHITE
    muted = samco_template.COLOR_SUBTITLE
    title_slide = deck.slides.add_slide(deck.slide_layouts[6])
    background = title_slide.background.fill
    background.solid(); background.fore_color.rgb = navy
    title_box = title_slide.shapes.add_textbox(Inches(0.55), Inches(2.2), Inches(12.2), Inches(1.1))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "SAMCO-PCO Project Controls Report"
    title_p.font.size = Pt(38); title_p.font.bold = True; title_p.font.color.rgb = white
    subtitle = title_slide.shapes.add_textbox(Inches(0.58), Inches(3.45), Inches(12.0), Inches(0.8))
    subtitle_p = subtitle.text_frame.paragraphs[0]
    subtitle_p.text = f"{project.get('project_display_name') or project.get('project_id')} | 30 data-driven control sections | {generated_at}"
    subtitle_p.font.size = Pt(16); subtitle_p.font.color.rgb = RGBColor(205, 224, 240)
    for item in sections:
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, deck.slide_width, Inches(0.82))
        header.fill.solid(); header.fill.fore_color.rgb = navy; header.line.fill.background()
        title = slide.shapes.add_textbox(Inches(0.32), Inches(0.14), Inches(10.5), Inches(0.5))
        title_p = title.text_frame.paragraphs[0]
        title_p.text = f"{item['code']} | {item['title']}"
        title_p.font.size = Pt(22); title_p.font.bold = True; title_p.font.color.rgb = white
        state = slide.shapes.add_textbox(Inches(9.3), Inches(1.0), Inches(3.55), Inches(0.35))
        state_p = state.text_frame.paragraphs[0]
        state_p.text = item["status"]; state_p.font.size = Pt(11); state_p.font.bold = True; state_p.font.color.rgb = RGBColor(10, 125, 69) if item["rows"] else RGBColor(170, 107, 0)
        sources = slide.shapes.add_textbox(Inches(0.4), Inches(1.02), Inches(8.6), Inches(1.0))
        source_frame = sources.text_frame; source_frame.word_wrap = True
        source_frame.paragraphs[0].text = "CSV sources"
        source_frame.paragraphs[0].font.size = Pt(12); source_frame.paragraphs[0].font.bold = True; source_frame.paragraphs[0].font.color.rgb = navy
        for source in item["sources"] or [{"logical_table": "No mapped source", "path": "", "row_count": 0}]:
            paragraph = source_frame.add_paragraph(); paragraph.text = f"{source['logical_table']}: {source['path']} ({source['row_count']} rows)"; paragraph.font.size = Pt(10); paragraph.font.color.rgb = muted
        if item["rows"]:
            headers = list(item["rows"][0].keys())[:6]
            table_shape = slide.shapes.add_table(min(13, len(item["rows"]) + 1), len(headers), Inches(0.4), Inches(2.15), Inches(12.45), Inches(4.6))
            table = table_shape.table
            for col, header_value in enumerate(headers):
                cell = table.cell(0, col); cell.text = header_value; cell.fill.solid(); cell.fill.fore_color.rgb = navy
                cell.text_frame.paragraphs[0].font.size = Pt(8); cell.text_frame.paragraphs[0].font.bold = True; cell.text_frame.paragraphs[0].font.color.rgb = white
            for row_index, row in enumerate(item["rows"][:12], start=1):
                for col, header_value in enumerate(headers):
                    cell = table.cell(row_index, col); cell.text = str(row.get(header_value, "") or "")[:80]
                    cell.text_frame.paragraphs[0].font.size = Pt(7)
        else:
            note = slide.shapes.add_textbox(Inches(0.4), Inches(2.25), Inches(12.0), Inches(0.6))
            note_p = note.text_frame.paragraphs[0]
            note_p.text = "No controlled CSV data is available for this report section. No value or conclusion has been inferred."
            note_p.font.size = Pt(16); note_p.font.color.rgb = muted
    deck.save(path)


def ensure_samco_pco_project_report(project: dict[str, Any], output_dir: Path, public_slug: str) -> dict[str, Any]:
    """Generate an HTML/PDF/PPTX report triplet from selected-project CSV data."""
    output_dir.mkdir(parents=True, exist_ok=True)
    generated_at = datetime.now(timezone.utc).replace(microsecond=0).isoformat()
    sections = _section_data(project)
    html_path = output_dir / f"{REPORT_STEM}.html"
    pdf_path = output_dir / f"{REPORT_STEM}.pdf"
    pptx_path = output_dir / f"{REPORT_STEM}.pptx"
    _write_html(html_path, project, sections, generated_at)
    _write_pdf(pdf_path, project, sections, generated_at)
    _write_pptx(pptx_path, project, sections, generated_at)
    return {
        "html": f"/generated/{public_slug}/{html_path.name}",
        "pdf": f"/generated/{public_slug}/{pdf_path.name}",
        "pptx": f"/generated/{public_slug}/{pptx_path.name}",
        "files": {
            extension: {"name": path.name, "bytes": path.stat().st_size, "sha256": _sha256(path)}
            for extension, path in (("html", html_path), ("pdf", pdf_path), ("pptx", pptx_path))
        },
        "source_project_id": project.get("project_id"),
        "source_report_fingerprint": project.get("fingerprint"),
        "generator": {
            "name": "SAMCO-PCO Report Generator",
            "mode": "data-driven project CSV adapter",
            "template_sha256": _sha256(ENGINE_SOURCE),
            "version": REPORT_GENERATOR_VERSION,
        },
    }
