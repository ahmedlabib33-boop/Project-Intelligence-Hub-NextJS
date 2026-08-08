from __future__ import annotations

"""Create project-scoped Output Studio artifacts from canonical project data.

The website only serves static assets.  This module therefore produces the
downloadable PDF and PowerPoint counterparts locally, alongside the existing
HTML reports, before the Vercel build is published.
"""

import hashlib
import html
import json
import os
import shutil
import subprocess
import sys
from datetime import datetime
from pathlib import Path
from typing import Any


REPORTS: tuple[tuple[str, str, str], ...] = (
    ("executive_dashboard", "01_executive_dashboard", "Executive Dashboard"),
    ("master_dashboard", "02_master_dashboard", "Master Dashboard"),
    ("elite_svg_charts", "03_elite_svg_charts", "Elite SVG Charts"),
    ("linked_executive_dashboard", "04_linked_executive_dashboard", "Linked Executive Dashboard"),
)
REPORT_GENERATOR_VERSION = "2026.08.project-scoped-html-pdf-pptx.v6-source-html"
WORKSPACE_ROOT = Path(__file__).resolve().parents[1]
CANONICAL_OUTPUTS_ROOT = WORKSPACE_ROOT / "11-outputs"


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _text(value: Any) -> str:
    if value is None:
        return "N/A"
    return str(value)


def _money(value: Any) -> str:
    if isinstance(value, (int, float)):
        return f"EGP {value:,.0f}"
    return "N/A"


def _percent(value: Any) -> str:
    if isinstance(value, (int, float)):
        return f"{value * 100:.1f}%"
    return "N/A"


def _metric_rows(project: dict[str, Any]) -> list[tuple[str, str]]:
    return [
        ("Project", _text(project.get("project_display_name"))),
        ("Project ID", _text(project.get("project_id"))),
        ("Sector", _text(project.get("sector"))),
        ("Status", _text(project.get("status"))),
        ("Contract Value", _money(project.get("contract_value"))),
        ("Paid Amount", _money(project.get("paid_amount"))),
        ("Actual Progress", _percent(project.get("actual_progress"))),
        ("Planned Progress", _percent(project.get("planned_progress"))),
        ("BAC / PV / EV / AC", " / ".join(_money(project.get(key)) for key in ("bac", "pv", "ev", "ac"))),
        ("SPI / CPI", " / ".join(_text(round(project[key], 2)) if isinstance(project.get(key), (int, float)) else "N/A" for key in ("spi", "cpi"))),
        ("Delay Days", _text(project.get("delay_days"))),
        ("Last Source Update", _text(project.get("last_updated"))),
    ]


def _write_fallback_html(path: Path, title: str, project: dict[str, Any]) -> None:
    """Create a controlled fallback when no approved source HTML exists.

    The fallback is deliberately not presented as a replacement for an approved
    report.  It keeps the selected-project output usable without publishing a
    report that belongs to a different project.
    """
    metrics = "".join(
        f"<tr><th>{html.escape(label)}</th><td>{html.escape(value)}</td></tr>"
        for label, value in _metric_rows(project)
    )
    path.write_text(
        "<!doctype html><html><head><meta charset='utf-8'>"
        f"<title>{html.escape(title)} - {html.escape(_text(project.get('project_display_name')))}</title>"
        "<style>body{font-family:Arial,sans-serif;margin:32px;color:#0a2340}h1{color:#073b66}"
        "table{border-collapse:collapse;width:100%;max-width:1100px}th,td{padding:10px;border:1px solid #cbd5e1;text-align:left}"
        "th{background:#e2e8f0;width:32%}small{color:#475569}</style></head><body>"
        f"<h1>{html.escape(title)}</h1><p><b>{html.escape(_text(project.get('project_display_name')))}</b></p>"
        f"<table>{metrics}</table><p><small>Generated from the selected project only.</small></p></body></html>",
        encoding="utf-8",
    )


def _canonical_html_sources(project: dict[str, Any]) -> tuple[dict[str, Path], dict[str, Any]]:
    """Return only the legacy HTML reports proven to belong to this project."""
    folder_name = str(project.get("project_folder_name") or "").strip()
    project_id = str(project.get("project_id") or "").strip()
    if not folder_name or not project_id:
        return {}, {}
    source_dir = CANONICAL_OUTPUTS_ROOT / folder_name
    manifest_path = source_dir / ".output_manifest.json"
    try:
        source_manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    except (OSError, UnicodeDecodeError, json.JSONDecodeError):
        return {}, {}
    if str(source_manifest.get("project_id") or "").strip() != project_id:
        return {}, {}
    sources: dict[str, Path] = {}
    for key, stem, _ in REPORTS:
        candidate = source_dir / f"{stem}.html"
        if candidate.exists() and candidate.stat().st_size > 4096:
            sources[key] = candidate
    return sources, source_manifest


def _publish_source_html_or_fallback(
    html_path: Path,
    key: str,
    title: str,
    project: dict[str, Any],
    source_reports: dict[str, Path],
) -> str:
    """Publish an approved project-owned source report, never a shared report."""
    source = source_reports.get(key)
    if source is not None:
        # In the self-contained D: workspace the approved report can already
        # be the target file.  Avoid SameFileError while retaining it intact.
        if source.resolve() != html_path.resolve():
            shutil.copy2(source, html_path)
        return "canonical_project_html"
    _write_fallback_html(html_path, title, project)
    return "controlled_fallback"


def _chrome_path() -> str | None:
    candidates = [
        shutil.which("chrome"),
        shutil.which("msedge"),
        r"C:\\Program Files\\Google\\Chrome\\Application\\chrome.exe",
        r"C:\\Program Files (x86)\\Google\\Chrome\\Application\\chrome.exe",
        r"C:\\Program Files (x86)\\Microsoft\\Edge\\Application\\msedge.exe",
    ]
    return next((candidate for candidate in candidates if candidate and Path(candidate).exists()), None)


def _write_pdf(path: Path, title: str, project: dict[str, Any], html_path: Path) -> None:
    """Print the approved HTML first; retain a deterministic A3 fallback."""
    chrome = _chrome_path()
    if chrome:
        try:
            result = subprocess.run(
                [
                    chrome,
                    "--headless=new",
                    "--disable-gpu",
                    "--no-pdf-header-footer",
                    f"--print-to-pdf={path}",
                    html_path.resolve().as_uri(),
                ],
                check=False,
                capture_output=True,
                text=True,
                timeout=90,
            )
            if result.returncode == 0 and path.exists() and path.stat().st_size > 0:
                return
        except (OSError, subprocess.SubprocessError):
            pass
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A3, landscape
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
    except ImportError as exc:  # validated by the pipeline before publishing
        raise RuntimeError("Output Studio PDF export requires reportlab.") from exc

    styles = getSampleStyleSheet()
    document = SimpleDocTemplate(
        str(path), pagesize=landscape(A3), leftMargin=15 * mm, rightMargin=15 * mm, topMargin=15 * mm, bottomMargin=15 * mm
    )
    data = [[Paragraph("<b>Metric</b>", styles["BodyText"]), Paragraph("<b>Value</b>", styles["BodyText"])] ]
    for label, value in _metric_rows(project):
        data.append([Paragraph(html.escape(label), styles["BodyText"]), Paragraph(html.escape(value), styles["BodyText"])])
    table = Table(data, colWidths=[95 * mm, 270 * mm], repeatRows=1)
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#dce6f1")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#092c4c")),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#9fb3c8")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 9),
        ("RIGHTPADDING", (0, 0), (-1, -1), 9),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
    ]))
    document.build([
        Paragraph(title, styles["Title"]),
        Paragraph(html.escape(_text(project.get("project_display_name"))), styles["Heading2"]),
        Spacer(1, 8 * mm),
        table,
        Spacer(1, 6 * mm),
        Paragraph("Source-controlled project report. Delay and EOT conclusions remain indicative until Primavera P6 recalculation is verified.", styles["BodyText"]),
    ])


def _write_pptx(path: Path, title: str, project: dict[str, Any]) -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
    except ImportError as exc:  # validated by the pipeline before publishing
        raise RuntimeError("Output Studio PowerPoint export requires python-pptx.") from exc

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    navy = RGBColor(8, 37, 66)
    cyan = RGBColor(57, 215, 210)
    white = RGBColor(244, 250, 255)
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = navy
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.4), Inches(12.2), Inches(0.75))
    paragraph = title_box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(30)
    paragraph.font.bold = True
    paragraph.font.color.rgb = white
    subtitle = slide.shapes.add_textbox(Inches(0.58), Inches(1.12), Inches(12.0), Inches(0.4))
    subtitle_p = subtitle.text_frame.paragraphs[0]
    subtitle_p.text = f"{_text(project.get('project_display_name'))} | {_text(project.get('project_id'))}"
    subtitle_p.font.size = Pt(14)
    subtitle_p.font.color.rgb = cyan
    metrics = _metric_rows(project)[:10]
    for index, (label, value) in enumerate(metrics):
        row, column = divmod(index, 2)
        left = Inches(0.6 + column * 6.25)
        top = Inches(1.75 + row * 1.02)
        tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.8), Inches(0.82))
        tile.fill.solid()
        tile.fill.fore_color.rgb = RGBColor(13, 61, 99)
        tile.line.color.rgb = cyan
        frame = tile.text_frame
        frame.clear()
        label_p = frame.paragraphs[0]
        label_p.text = label.upper()
        label_p.font.size = Pt(8)
        label_p.font.bold = True
        label_p.font.color.rgb = cyan
        value_p = frame.add_paragraph()
        value_p.text = value
        value_p.font.size = Pt(13)
        value_p.font.bold = True
        value_p.font.color.rgb = white
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.0), Inches(0.25))
    footer_p = footer.text_frame.paragraphs[0]
    footer_p.text = f"Generated {datetime.now().strftime('%d %b %Y')} | Source-controlled selected-project output"
    footer_p.font.size = Pt(8)
    footer_p.font.color.rgb = RGBColor(167, 195, 214)
    presentation.save(path)


def _ensure_controlled_tia_artifacts(project: dict[str, Any], output_dir: Path, public_slug: str) -> dict[str, Any] | None:
    """Publish a TIA report only from an approved controlled project run.

    Historic generic and submitted-guide TIA outputs are deliberately excluded
    from new Output Studio manifests.  A controlled run that is setup,
    conditional, or awaiting reconciliation is evidence review material, not a
    publishable EOT report.
    """
    features = project.get("features") if isinstance(project.get("features"), dict) else {}
    delay = features.get("delay_analysis") if isinstance(features.get("delay_analysis"), dict) else {}
    controlled = delay.get("controlled_tia") if isinstance(delay.get("controlled_tia"), dict) else {}
    if controlled.get("status") != "READY_AND_CALCULATED" or controlled.get("approval_status") != "approved":
        return None
    # A later controlled-report template can consume this same run.  Until a
    # project passes all gates, it is safer not to emit a client-facing report.
    return None


def ensure_project_report_artifacts(
    project: dict[str, Any], output_dir: Path, public_slug: str | None = None
) -> dict[str, dict[str, Any]]:
    """Create or refresh the HTML/PDF/PPTX triplet for every primary report."""
    output_dir.mkdir(parents=True, exist_ok=True)
    project_slug = str(public_slug or project.get("project_key") or project.get("project_id") or "project")
    manifest_path = output_dir / ".report_manifest.json"
    if manifest_path.exists():
        try:
            existing = json.loads(manifest_path.read_text(encoding="utf-8"))
            existing_reports = existing.get("reports", {})
            if (
                existing.get("project_fingerprint") == project.get("fingerprint")
                and existing.get("generator_version") == REPORT_GENERATOR_VERSION
                and all(
                    (output_dir / f"{stem}.{extension}").exists()
                    for _, stem, _ in REPORTS
                    for extension in ("html", "pdf", "pptx")
                )
                and existing_reports
            ):
                return existing_reports
        except Exception:
            pass
    source_reports, source_manifest = _canonical_html_sources(project)
    results: dict[str, dict[str, Any]] = {}
    for key, stem, title in REPORTS:
        html_path = output_dir / f"{stem}.html"
        pdf_path = output_dir / f"{stem}.pdf"
        pptx_path = output_dir / f"{stem}.pptx"
        html_origin = _publish_source_html_or_fallback(html_path, key, title, project, source_reports)
        _write_pdf(pdf_path, title, project, html_path)
        _write_pptx(pptx_path, title, project)
        results[key] = {
            "html": f"/generated/{project_slug}/{html_path.name}",
            "pdf": f"/generated/{project_slug}/{pdf_path.name}",
            "pptx": f"/generated/{project_slug}/{pptx_path.name}",
            "files": {
                extension: {"name": path.name, "bytes": path.stat().st_size, "sha256": _sha256(path)}
                for extension, path in (("html", html_path), ("pdf", pdf_path), ("pptx", pptx_path))
            },
            "html_origin": html_origin,
            "source_project_id": source_manifest.get("project_id") if html_origin == "canonical_project_html" else project.get("project_id"),
            "source_report_fingerprint": source_manifest.get("fingerprint") if html_origin == "canonical_project_html" else project.get("fingerprint"),
        }
    tia_artifacts = _ensure_controlled_tia_artifacts(project, output_dir, project_slug)
    if tia_artifacts:
        results["tia_controlled_assessment"] = tia_artifacts
    feature_payload = project.get("features") if isinstance(project.get("features"), dict) else {}
    assessment = feature_payload.get("four_pipeline") if isinstance(feature_payload.get("four_pipeline"), dict) else {}
    manifest = {
        "project_id": project.get("project_id"),
        "project_key": project_slug,
        "project_fingerprint": project.get("fingerprint"),
        "generator_version": REPORT_GENERATOR_VERSION,
        "chart_catalog_version": (
            project.get("chart_payloads", {}).get("catalog_version")
            if isinstance(project.get("chart_payloads"), dict)
            else None
        ),
        "chart_status": {
            str(item.get("id")): str(item.get("status"))
            for item in (project.get("chart_payloads", {}).get("charts", []) if isinstance(project.get("chart_payloads"), dict) else [])
            if isinstance(item, dict) and item.get("id")
        },
        "assessment": {
            "analysis_run_id": assessment.get("analysis_run_id"),
            "assessment_profile": assessment.get("assessment_profile"),
            "assessment_status": assessment.get("assessment_status"),
            "determination_status": assessment.get("determination_status"),
            "source_scope": assessment.get("source_scope"),
        },
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "reports": results,
    }
    (output_dir / ".report_manifest.json").write_text(json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8")
    return results
